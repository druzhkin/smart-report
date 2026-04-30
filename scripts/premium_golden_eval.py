"""Evaluate a generated report package against the 10k RUB product bar.

This is a local, deterministic harness for golden tasks. It does not call live
LLMs or research providers. Feed it the exported DOCX/PPTX plus optional audit
and artifact-QA JSON files; it returns a comparable scorecard.
"""

from __future__ import annotations

import argparse
import csv
import json
import tempfile
import zipfile
from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Any


@dataclass
class GoldenEvalResult:
    label: str
    overall_score: int
    verdict: str
    subscores: dict[str, int]
    metrics: dict[str, Any]
    blockers: list[str] = field(default_factory=list)
    recommendations: list[str] = field(default_factory=list)


def evaluate_package(
    *,
    label: str,
    docx_path: Path | None = None,
    pptx_path: Path | None = None,
    audit_json_path: Path | None = None,
    artifact_qa_json_path: Path | None = None,
) -> GoldenEvalResult:
    metrics: dict[str, Any] = {}
    blockers: list[str] = []
    recommendations: list[str] = []

    docx_metrics = _docx_metrics(docx_path) if docx_path else {}
    pptx_metrics = _pptx_metrics(pptx_path) if pptx_path else {}
    audit = _read_json(audit_json_path) if audit_json_path else {}
    artifact_qa = _read_json(artifact_qa_json_path) if artifact_qa_json_path else {}

    metrics.update({f"docx_{key}": value for key, value in docx_metrics.items()})
    metrics.update({f"pptx_{key}": value for key, value in pptx_metrics.items()})
    metrics["audit_present"] = bool(audit)
    metrics["artifact_qa_present"] = bool(artifact_qa)

    client_readiness = audit.get("client_readiness") or {}
    premium_readiness = audit.get("premium_readiness") or {}
    analytic_closure = audit.get("analytic_closure") or {}
    evidence_audit = audit.get("evidence_audit") or {}
    adjudication_audit = audit.get("adjudication_audit") or {}
    visual_review = audit.get("visual_review") or {}

    if premium_readiness:
        metrics["premium_score"] = int(premium_readiness.get("score") or 0)
        metrics["premium_ready"] = bool(premium_readiness.get("ready"))
    if client_readiness:
        metrics["client_ready"] = bool(client_readiness.get("ready"))
        metrics["client_score"] = int(client_readiness.get("score") or 0)
    if analytic_closure:
        metrics["closure_score"] = int(analytic_closure.get("overall_score") or 0)
        metrics["closure_open"] = int(analytic_closure.get("not_closed") or 0) + int(
            analytic_closure.get("not_started") or 0
        )
    if evidence_audit:
        metrics["evidence_support_score"] = int(evidence_audit.get("overall_score") or 0)
        metrics["evidence_unsupported_claims"] = int(evidence_audit.get("unsupported") or 0)
    if adjudication_audit:
        metrics["adjudication_score"] = int(adjudication_audit.get("overall_score") or 0)
        metrics["adjudication_unresolved"] = int(adjudication_audit.get("unresolved") or 0)
        metrics["adjudication_critical_unresolved"] = int(adjudication_audit.get("critical_unresolved") or 0)
    if visual_review:
        metrics["visual_review_status"] = visual_review.get("status")
        metrics["visual_review_ready"] = bool(visual_review.get("ready"))
    if artifact_qa:
        metrics["artifact_qa_status"] = artifact_qa.get("status")
        metrics["artifact_qa_issues"] = int((artifact_qa.get("summary") or {}).get("issues") or 0)

    subscores = {
        "content_depth": _content_depth_score(docx_metrics),
        "evidence": _evidence_score(audit),
        "claim_support": _claim_support_score(evidence_audit),
        "adjudication": _adjudication_score(adjudication_audit),
        "analytic_closure": _closure_score(analytic_closure),
        "design_package": _design_package_score(docx_metrics, pptx_metrics, artifact_qa),
        "visual_review": _visual_review_score(visual_review),
        "delivery_safety": _delivery_safety_score(client_readiness, premium_readiness, artifact_qa),
    }

    _collect_blockers(metrics, blockers, recommendations)
    overall = round(sum(subscores.values()) / len(subscores))
    if blockers:
        overall = min(overall, 74)
    verdict = _verdict(overall, blockers)
    return GoldenEvalResult(
        label=label,
        overall_score=overall,
        verdict=verdict,
        subscores=subscores,
        metrics=metrics,
        blockers=blockers,
        recommendations=recommendations,
    )


def evaluate_package_zip(*, label: str, package_zip_path: Path) -> GoldenEvalResult:
    if not package_zip_path.exists():
        return GoldenEvalResult(
            label=label,
            overall_score=0,
            verdict="not_acceptable",
            subscores={},
            metrics={"package_zip_exists": False},
            blockers=["Premium package ZIP does not exist."],
            recommendations=["Generate a premium package before evaluation."],
        )
    with tempfile.TemporaryDirectory(prefix="premium-golden-") as tmp:
        tmp_path = Path(tmp)
        with zipfile.ZipFile(package_zip_path) as zf:
            zf.extractall(tmp_path)
        return evaluate_package(
            label=label,
            docx_path=tmp_path / "01_premium_report.docx",
            pptx_path=tmp_path / "02_premium_deck.pptx",
            audit_json_path=tmp_path / "05_audit.json",
            artifact_qa_json_path=tmp_path / "07_artifact_qa.json",
        )


def _docx_metrics(path: Path | None) -> dict[str, Any]:
    if not path or not path.exists():
        return {"exists": False}
    try:
        from docx import Document

        doc = Document(path)
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        table_cells = [
            cell.text.strip()
            for table in doc.tables
            for row in table.rows
            for cell in row.cells
            if cell.text.strip()
        ]
        text = "\n".join([*paragraphs, *table_cells])
        return {
            "exists": True,
            "paragraphs": len(paragraphs),
            "tables": len(doc.tables),
            "text_chars": len(text),
            "estimated_pages": _estimate_docx_pages(text, len(doc.tables)),
            "page_sections": len(doc.sections),
            "url_count": text.count("http://") + text.count("https://"),
            "has_readiness_gate": "Premium Readiness Gate" in text,
            "has_scorecard": "Executive Evidence Scorecard" in text,
        }
    except Exception as exc:  # pragma: no cover - defensive CLI path
        return {"exists": True, "error": str(exc)}


def _estimate_docx_pages(text: str, table_count: int) -> int:
    return max(1, round(len(text) / 2800 + table_count * 0.15))


def _pptx_metrics(path: Path | None) -> dict[str, Any]:
    if not path or not path.exists():
        return {"exists": False}
    try:
        from pptx import Presentation

        deck = Presentation(str(path))
        text_parts = []
        tables = 0
        for slide in deck.slides:
            for shape in slide.shapes:
                if getattr(shape, "has_table", False):
                    tables += 1
                if hasattr(shape, "text") and shape.text:
                    text_parts.append(shape.text)
        text = "\n".join(text_parts)
        return {
            "exists": True,
            "slides": len(deck.slides),
            "tables": tables,
            "text_chars": len(text),
            "has_readiness": "Paid-Delivery Readiness" in text,
        }
    except Exception as exc:  # pragma: no cover - defensive CLI path
        return {"exists": True, "error": str(exc)}


def _content_depth_score(docx: dict[str, Any]) -> int:
    if not docx.get("exists"):
        return 0
    score = 0
    score += min(35, int(docx.get("text_chars", 0) / 500))
    score += min(25, int(docx.get("tables", 0) * 2))
    score += 15 if docx.get("has_scorecard") else 0
    score += 15 if docx.get("has_readiness_gate") else 0
    score += min(10, int(docx.get("url_count", 0)))
    return min(100, score)


def _evidence_score(audit: dict[str, Any]) -> int:
    premium = audit.get("premium_readiness") or {}
    analysis = audit.get("analysis") or {}
    facts = analysis.get("high_relevance_facts") or analysis.get("all_numeric_facts") or []
    sources = ((audit.get("client_report") or {}).get("all_sources") or [])
    score = int(premium.get("score") or 0)
    score = max(score, min(100, len(facts) * 8 + len(sources) * 4))
    return min(100, score)


def _closure_score(closure: dict[str, Any]) -> int:
    if not closure:
        return 0
    return min(100, int(closure.get("overall_score") or 0))


def _claim_support_score(evidence_audit: dict[str, Any]) -> int:
    if not evidence_audit:
        return 0
    score = min(100, int(evidence_audit.get("overall_score") or 0))
    unsupported = int(evidence_audit.get("unsupported") or 0)
    return max(0, score - unsupported * 10)


def _adjudication_score(adjudication_audit: dict[str, Any]) -> int:
    if not adjudication_audit:
        return 0
    score = min(100, int(adjudication_audit.get("overall_score") or 0))
    unresolved = int(adjudication_audit.get("unresolved") or 0)
    critical_unresolved = int(adjudication_audit.get("critical_unresolved") or 0)
    return max(0, score - unresolved * 8 - critical_unresolved * 20)


def _design_package_score(
    docx: dict[str, Any],
    pptx: dict[str, Any],
    artifact_qa: dict[str, Any],
) -> int:
    score = 0
    score += 25 if docx.get("exists") else 0
    score += 25 if pptx.get("exists") else 0
    score += min(20, int(pptx.get("slides", 0) * 2))
    score += min(15, int(docx.get("tables", 0)))
    status = artifact_qa.get("status")
    if status == "passed":
        score += 15
    elif status == "blocked":
        score += 5
    return min(100, score)


def _visual_review_score(visual_review: dict[str, Any]) -> int:
    if not visual_review:
        return 0
    if visual_review.get("ready"):
        return 100
    if visual_review.get("status") == "pending":
        return 35
    return 0


def _delivery_safety_score(
    client_readiness: dict[str, Any],
    premium_readiness: dict[str, Any],
    artifact_qa: dict[str, Any],
) -> int:
    score = 0
    score += 35 if client_readiness.get("ready") else 0
    score += 40 if premium_readiness.get("ready") else min(25, int(premium_readiness.get("score") or 0) // 4)
    score += 25 if artifact_qa.get("status") == "passed" else 0
    return min(100, score)


def _collect_blockers(
    metrics: dict[str, Any],
    blockers: list[str],
    recommendations: list[str],
) -> None:
    if not metrics.get("audit_present"):
        blockers.append("Audit/readiness JSON is missing.")
        recommendations.append("Evaluate paid delivery only with audit, readiness, and closure metadata.")
    if metrics.get("artifact_qa_status") != "passed":
        blockers.append("Visual artifact QA has not passed.")
        recommendations.append("Run LibreOffice/Poppler rendering and inspect generated pages/slides.")
    if not metrics.get("premium_ready"):
        blockers.append("Premium readiness gate is not ready.")
        recommendations.append("Close premium readiness issues before client delivery.")
    if not metrics.get("client_ready"):
        blockers.append("Client readiness gate is not ready.")
        recommendations.append("Export only as draft until client readiness is green.")
    if int(metrics.get("closure_open") or 0) > 0:
        blockers.append("Analytic closure still has open leads.")
        recommendations.append("Run targeted follow-up or explicitly bracket limitations.")
    if int(metrics.get("evidence_unsupported_claims") or 0) > 0:
        blockers.append("Evidence audit has unsupported client-facing conclusions.")
        recommendations.append("Add citations, source links, or numeric fact backing to unsupported conclusions.")
    if int(metrics.get("adjudication_critical_unresolved") or 0) > 0:
        blockers.append("Adjudication audit has unresolved critical conflicts.")
        recommendations.append("Resolve critical conflicts or explicitly bracket them before paid delivery.")
    elif int(metrics.get("adjudication_unresolved") or 0) > 0:
        blockers.append("Adjudication audit has unresolved conflicts.")
        recommendations.append("Add resolution logic, scope boundaries, or limitations for unresolved conflicts.")
    if metrics.get("visual_review_status") and not metrics.get("visual_review_ready"):
        blockers.append("Manual visual review is not approved.")
        recommendations.append("Inspect rendered pages/slides and approve visual_review before paid delivery.")


def _verdict(score: int, blockers: list[str]) -> str:
    if score >= 85 and not blockers:
        return "paid_client_ready"
    if score >= 70:
        return "strong_draft_not_paid_ready" if blockers else "borderline_paid_ready"
    if score >= 50:
        return "useful_internal_draft"
    return "not_acceptable"


def _read_json(path: Path | None) -> dict[str, Any]:
    if not path or not path.exists():
        return {}
    return json.loads(path.read_text(encoding="utf-8"))


def main() -> int:
    parser = argparse.ArgumentParser(description="Evaluate a premium report package.")
    parser.add_argument(
        "--manifest",
        type=Path,
        help="JSON file with a list of golden tasks to evaluate as a leaderboard.",
    )
    parser.add_argument("--label", default="golden-task")
    parser.add_argument("--docx", type=Path)
    parser.add_argument("--pptx", type=Path)
    parser.add_argument("--audit-json", type=Path)
    parser.add_argument("--artifact-qa-json", type=Path)
    parser.add_argument("--package-zip", type=Path)
    parser.add_argument("--json", type=Path)
    parser.add_argument("--csv", type=Path)
    args = parser.parse_args()

    if args.manifest:
        payload = evaluate_manifest(args.manifest)
        print(json.dumps(payload, ensure_ascii=False, indent=2))
        if args.json:
            args.json.parent.mkdir(parents=True, exist_ok=True)
            args.json.write_text(json.dumps(payload, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
        if args.csv:
            _write_leaderboard_csv(args.csv, payload["results"])
        return 0

    if args.package_zip:
        result = evaluate_package_zip(label=args.label, package_zip_path=args.package_zip)
    else:
        result = evaluate_package(
            label=args.label,
            docx_path=args.docx,
            pptx_path=args.pptx,
            audit_json_path=args.audit_json,
            artifact_qa_json_path=args.artifact_qa_json,
        )
    payload = asdict(result)
    print(json.dumps(payload, ensure_ascii=False, indent=2))
    if args.json:
        args.json.parent.mkdir(parents=True, exist_ok=True)
        args.json.write_text(json.dumps(payload, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    if args.csv:
        args.csv.parent.mkdir(parents=True, exist_ok=True)
        _write_leaderboard_csv(args.csv, [payload])
    return 0


def evaluate_manifest(path: Path) -> dict[str, Any]:
    manifest = _read_json(path)
    tasks = manifest.get("tasks") if isinstance(manifest, dict) else manifest
    if not isinstance(tasks, list):
        raise SystemExit("Manifest must be a JSON list or an object with a 'tasks' list.")

    base_dir = path.parent
    results = []
    for idx, task in enumerate(tasks, start=1):
        if not isinstance(task, dict):
            raise SystemExit(f"Manifest task #{idx} must be an object.")
        label = str(task.get("label") or f"golden-{idx}")
        result = evaluate_package(
            label=label,
            docx_path=_optional_path(base_dir, task.get("docx")),
            pptx_path=_optional_path(base_dir, task.get("pptx")),
            audit_json_path=_optional_path(base_dir, task.get("audit_json")),
            artifact_qa_json_path=_optional_path(base_dir, task.get("artifact_qa_json")),
        ) if not task.get("package_zip") else evaluate_package_zip(
            label=label,
            package_zip_path=_optional_path(base_dir, task.get("package_zip")) or Path(),
        )
        results.append(asdict(result))

    scores = [int(item["overall_score"]) for item in results]
    paid_ready = sum(1 for item in results if item["verdict"] == "paid_client_ready")
    summary = {
        "tasks": len(results),
        "average_score": round(sum(scores) / len(scores), 1) if scores else 0,
        "minimum_score": min(scores) if scores else 0,
        "paid_client_ready": paid_ready,
        "paid_client_ready_rate": round(paid_ready / len(results), 3) if results else 0,
    }
    return {"summary": summary, "results": results}


def _optional_path(base_dir: Path, raw: Any) -> Path | None:
    if raw in (None, ""):
        return None
    path = Path(str(raw))
    return path if path.is_absolute() else base_dir / path


def _write_leaderboard_csv(path: Path, rows: list[dict[str, Any]]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    subscore_keys = sorted({key for row in rows for key in (row.get("subscores") or {}).keys()})
    metric_keys = sorted({key for row in rows for key in (row.get("metrics") or {}).keys()})
    with path.open("w", encoding="utf-8-sig", newline="") as fh:
        writer = csv.writer(fh)
        writer.writerow(
            [
                "label",
                "overall_score",
                "verdict",
                "blocker_count",
                *[f"subscore_{key}" for key in subscore_keys],
                *[f"metric_{key}" for key in metric_keys],
            ]
        )
        for row in sorted(rows, key=lambda item: int(item.get("overall_score") or 0), reverse=True):
            writer.writerow(
                [
                    row.get("label"),
                    row.get("overall_score"),
                    row.get("verdict"),
                    len(row.get("blockers") or []),
                    *[(row.get("subscores") or {}).get(key, "") for key in subscore_keys],
                    *[(row.get("metrics") or {}).get(key, "") for key in metric_keys],
                ]
            )


if __name__ == "__main__":
    raise SystemExit(main())
