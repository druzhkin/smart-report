"""QA checks for premium Smart Report DOCX and PPTX artifacts.

The checker is intentionally additive: it does not change exported files. It
performs structural inspection everywhere and performs visual rendering only
when LibreOffice (`soffice`) and Poppler (`pdftoppm`) are available locally.
"""

from __future__ import annotations

import argparse
import html
import json
import shutil
import subprocess
import sys
import tempfile
from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Any

INTERNAL_MARKERS = (
    "[STRONG]",
    "[MODERATE]",
    "[WEAK]",
    "[SPECULATIVE]",
    "[REF:",
    "main_synthesis",
    "consensus_section",
    "gaps_filled_section",
    "source_reports",
    "followup_reports",
)

VISUAL_REVIEW_RUBRIC = (
    "No text overflow, clipping, or overlapping elements.",
    "Tables are readable without broken rows, orphan headers, or cramped cells.",
    "Headings create a clear executive-to-detail hierarchy.",
    "Page and slide breaks preserve complete ideas and do not strand captions.",
    "Charts, scorecards, and badges are visually aligned and easy to scan.",
    "The package looks like a finished paid report, not a raw model export.",
)


@dataclass
class ArtifactQaResult:
    path: str
    kind: str
    exists: bool
    structural_status: str = "not_run"
    render_status: str = "not_run"
    missing_tools: list[str] = field(default_factory=list)
    metrics: dict[str, Any] = field(default_factory=dict)
    issues: list[str] = field(default_factory=list)
    rendered_files: list[str] = field(default_factory=list)

    @property
    def passed(self) -> bool:
        return self.exists and self.structural_status == "passed" and not self.issues


def run_qa(
    *,
    docx_path: Path | None = None,
    pptx_path: Path | None = None,
    out_dir: Path,
    render: bool = True,
) -> dict[str, Any]:
    out_dir.mkdir(parents=True, exist_ok=True)
    results: list[ArtifactQaResult] = []
    if docx_path:
        results.append(_inspect_docx(docx_path))
    if pptx_path:
        results.append(_inspect_pptx(pptx_path))

    if render:
        soffice = _find_tool(
            "soffice",
            [
                Path("C:/Program Files/LibreOffice/program/soffice.exe"),
                Path("C:/Program Files (x86)/LibreOffice/program/soffice.exe"),
            ],
        )
        pdftoppm = _find_tool(
            "pdftoppm",
            list(Path.home().glob(
                "AppData/Local/Microsoft/WinGet/Packages/oschwartz10612.Poppler_*/poppler-*/Library/bin/pdftoppm.exe"
            )),
        )
        for result in results:
            _render_artifact(result, out_dir, soffice=soffice, pdftoppm=pdftoppm)
    else:
        for result in results:
            result.render_status = "skipped"

    report = {
        "status": _overall_status(results),
        "summary": {
            "artifacts": len(results),
            "passed_structural": sum(1 for item in results if item.structural_status == "passed"),
            "rendered": sum(1 for item in results if item.render_status == "passed"),
            "blocked_render": sum(1 for item in results if item.render_status == "blocked"),
            "issues": sum(len(item.issues) for item in results),
        },
        "results": [asdict(item) for item in results],
        "manual_visual_review_rubric": list(VISUAL_REVIEW_RUBRIC),
    }
    render_index = _write_render_index(results, out_dir)
    if render_index:
        report["render_index"] = str(render_index)
    return report


def _inspect_docx(path: Path) -> ArtifactQaResult:
    result = ArtifactQaResult(path=str(path), kind="docx", exists=path.exists())
    if not result.exists:
        result.structural_status = "failed"
        result.issues.append("DOCX file does not exist.")
        return result

    try:
        from docx import Document

        doc = Document(path)
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        table_text = [
            cell.text.strip()
            for table in doc.tables
            for row in table.rows
            for cell in row.cells
            if cell.text.strip()
        ]
        text = "\n".join([*paragraphs, *table_text])
        headings = [
            p.text.strip()
            for p in doc.paragraphs
            if p.style and p.style.name.startswith("Heading") and p.text.strip()
        ]
        result.metrics = {
            "paragraphs": len(paragraphs),
            "tables": len(doc.tables),
            "sections": len(doc.sections),
            "headings": len(headings),
            "text_chars": len(text),
            "estimated_pages": _estimate_docx_pages(text, len(doc.tables)),
            "has_cover_brand": "SMART REPORT | PREMIUM ANALYTICAL REPORT" in text,
            "has_decision_dashboard": "Client Decision Dashboard" in text,
            "has_scorecard": "Executive Evidence Scorecard" in text,
            "has_readiness_gate": "Premium Readiness Gate" in text,
            "has_report_structure": "Report Structure" in text,
        }
        _add_common_content_issues(result, text)
        _require_metric(result, "paragraphs", 25, "DOCX has too few paragraphs for a long-form report.")
        _require_metric(result, "tables", 4, "DOCX has too few visual tables.")
        _require_metric(result, "text_chars", 5000, "DOCX text volume is below structural QA minimum.")
        for key, message in {
            "has_cover_brand": "DOCX cover brand marker is missing.",
            "has_decision_dashboard": "DOCX client decision dashboard is missing.",
            "has_scorecard": "DOCX evidence scorecard is missing.",
            "has_readiness_gate": "DOCX readiness gate is missing.",
            "has_report_structure": "DOCX report structure section is missing.",
        }.items():
            if not result.metrics.get(key):
                result.issues.append(message)
        result.structural_status = "passed" if not result.issues else "failed"
    except Exception as exc:  # pragma: no cover - defensive CLI path
        result.structural_status = "failed"
        result.issues.append(f"DOCX inspection failed: {exc}")
    return result


def _estimate_docx_pages(text: str, table_count: int) -> int:
    """Approximate page count when render tools are unavailable.

    It is deliberately conservative and only a metric; real paid-delivery page
    QA still requires DOCX/PPTX rendering through LibreOffice/Poppler.
    """
    text_pages = len(text) / 2800
    table_pages = table_count * 0.15
    return max(1, round(text_pages + table_pages))


def _find_tool(name: str, candidates: list[Path]) -> str | None:
    found = shutil.which(name)
    if found:
        return found
    for candidate in candidates:
        if candidate.exists():
            return str(candidate)
    return None


def _inspect_pptx(path: Path) -> ArtifactQaResult:
    result = ArtifactQaResult(path=str(path), kind="pptx", exists=path.exists())
    if not result.exists:
        result.structural_status = "failed"
        result.issues.append("PPTX file does not exist.")
        return result

    try:
        from pptx import Presentation

        deck = Presentation(str(path))
        slide_texts: list[str] = []
        shape_count = 0
        table_count = 0
        for slide in deck.slides:
            for shape in slide.shapes:
                shape_count += 1
                if getattr(shape, "has_table", False):
                    table_count += 1
                if hasattr(shape, "text") and shape.text:
                    slide_texts.append(shape.text.strip())
        text = "\n".join(part for part in slide_texts if part)
        result.metrics = {
            "slides": len(deck.slides),
            "shapes": shape_count,
            "tables": table_count,
            "text_chars": len(text),
            "has_executive_answer": "Executive Answer" in text,
            "has_readiness": "Paid-Delivery Readiness" in text,
            "has_evidence_base": "Evidence Base" in text,
        }
        _add_common_content_issues(result, text)
        _require_metric(result, "slides", 10, "PPTX has too few slides for the premium deck.")
        _require_metric(result, "tables", 3, "PPTX has too few table-based analytical slides.")
        _require_metric(result, "text_chars", 1000, "PPTX text volume is below structural QA minimum.")
        for key, message in {
            "has_executive_answer": "PPTX executive answer slide is missing.",
            "has_readiness": "PPTX readiness slide is missing.",
            "has_evidence_base": "PPTX evidence base slide is missing.",
        }.items():
            if not result.metrics.get(key):
                result.issues.append(message)
        result.structural_status = "passed" if not result.issues else "failed"
    except Exception as exc:  # pragma: no cover - defensive CLI path
        result.structural_status = "failed"
        result.issues.append(f"PPTX inspection failed: {exc}")
    return result


def _render_artifact(
    result: ArtifactQaResult,
    out_dir: Path,
    *,
    soffice: str | None,
    pdftoppm: str | None,
) -> None:
    if not result.exists:
        result.render_status = "failed"
        return
    missing = []
    if not soffice:
        missing.append("soffice")
    if not pdftoppm:
        missing.append("pdftoppm")
    if missing:
        result.render_status = "blocked"
        result.missing_tools = missing
        result.issues.append(
            "Visual render QA was not completed because required tools are missing: "
            + ", ".join(missing)
            + "."
        )
        return

    artifact = Path(result.path)
    artifact_out = out_dir / artifact.stem
    artifact_out.mkdir(parents=True, exist_ok=True)
    profile = Path(tempfile.mkdtemp(prefix="smart-report-lo-"))
    try:
        subprocess.run(
            [
                soffice,
                f"-env:UserInstallation=file:///{profile.as_posix()}",
                "--headless",
                "--convert-to",
                "pdf",
                "--outdir",
                str(artifact_out),
                str(artifact),
            ],
            check=True,
            capture_output=True,
            text=True,
            timeout=120,
        )
        pdf_path = artifact_out / f"{artifact.stem}.pdf"
        if not pdf_path.exists():
            result.render_status = "failed"
            result.issues.append("LibreOffice conversion did not produce a PDF.")
            return
        prefix = artifact_out / artifact.stem
        subprocess.run(
            [pdftoppm, "-png", str(pdf_path), str(prefix)],
            check=True,
            capture_output=True,
            text=True,
            timeout=120,
        )
        rendered = sorted(str(item) for item in artifact_out.glob(f"{artifact.stem}-*.png"))
        result.rendered_files = [str(pdf_path), *rendered]
        result.render_status = "passed" if rendered else "failed"
        if rendered:
            if result.kind == "docx":
                result.metrics["rendered_pages"] = len(rendered)
            elif result.kind == "pptx":
                result.metrics["rendered_slides"] = len(rendered)
        if not rendered:
            result.issues.append("PDF rendering produced no PNG pages.")
    except subprocess.SubprocessError as exc:
        result.render_status = "failed"
        result.issues.append(f"Visual render command failed: {exc}")
    finally:
        shutil.rmtree(profile, ignore_errors=True)


def _write_render_index(results: list[ArtifactQaResult], out_dir: Path) -> Path | None:
    image_entries: list[tuple[str, str, Path]] = []
    for result in results:
        for rendered_file in result.rendered_files:
            rendered_path = Path(rendered_file)
            if rendered_path.suffix.lower() != ".png":
                continue
            try:
                rel = rendered_path.relative_to(out_dir)
            except ValueError:
                rel = rendered_path
            image_entries.append((result.kind.upper(), rendered_path.stem, rel))
    if not image_entries:
        return None

    cards = []
    for kind, label, rel in image_entries:
        href = html.escape(rel.as_posix())
        cards.append(
            "\n".join(
                [
                    '<article class="card">',
                    f'  <div class="meta">{html.escape(kind)} · {html.escape(label)}</div>',
                    f'  <a href="{href}"><img src="{href}" alt="{html.escape(label)}" /></a>',
                    "</article>",
                ]
            )
        )

    page = f"""<!doctype html>
<html lang="en">
<head>
  <meta charset="utf-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1" />
  <title>Smart Report Artifact QA</title>
  <style>
    :root {{
      color-scheme: light;
      --ink: #111827;
      --muted: #6b7280;
      --line: #d7dce5;
      --paper: #f6f7f9;
      --panel: #ffffff;
    }}
    * {{ box-sizing: border-box; }}
    body {{
      margin: 0;
      background: var(--paper);
      color: var(--ink);
      font: 14px/1.45 Arial, Helvetica, sans-serif;
    }}
    header {{
      padding: 24px 28px 18px;
      border-bottom: 1px solid var(--line);
      background: var(--panel);
    }}
    h1 {{
      margin: 0 0 6px;
      font-size: 22px;
      line-height: 1.2;
      letter-spacing: 0;
    }}
    p {{ margin: 0; color: var(--muted); }}
    .rubric {{
      display: grid;
      grid-template-columns: repeat(auto-fit, minmax(260px, 1fr));
      gap: 8px 18px;
      margin-top: 16px;
      padding: 14px 0 0;
      border-top: 1px solid var(--line);
      color: var(--ink);
    }}
    .rubric div {{
      display: flex;
      gap: 8px;
      align-items: flex-start;
      color: var(--muted);
      font-size: 13px;
    }}
    .rubric b {{
      color: var(--ink);
      font-weight: 700;
    }}
    main {{
      display: grid;
      grid-template-columns: repeat(auto-fit, minmax(280px, 1fr));
      gap: 18px;
      padding: 22px;
    }}
    .card {{
      background: var(--panel);
      border: 1px solid var(--line);
      border-radius: 8px;
      overflow: hidden;
      box-shadow: 0 1px 2px rgba(15, 23, 42, 0.08);
    }}
    .meta {{
      padding: 10px 12px;
      border-bottom: 1px solid var(--line);
      color: var(--muted);
      font-size: 12px;
      text-transform: uppercase;
      letter-spacing: .08em;
    }}
    img {{
      display: block;
      width: 100%;
      height: auto;
      background: white;
    }}
  </style>
</head>
<body>
  <header>
    <h1>Smart Report Artifact QA</h1>
    <p>Rendered DOCX pages and PPTX slides for fast human visual review.</p>
    <section class="rubric" aria-label="Manual visual review rubric">
      {"".join(f"<div><b>{idx}.</b><span>{html.escape(item)}</span></div>" for idx, item in enumerate(VISUAL_REVIEW_RUBRIC, start=1))}
    </section>
  </header>
  <main>
    {"".join(cards)}
  </main>
</body>
</html>
"""
    index_path = out_dir / "index.html"
    index_path.write_text(page, encoding="utf-8")
    return index_path


def _add_common_content_issues(result: ArtifactQaResult, text: str) -> None:
    for marker in INTERNAL_MARKERS:
        if marker in text:
            result.issues.append(f"Internal marker leaked into {result.kind.upper()}: {marker}")
    if "\ufffd" in text:
        result.issues.append(f"{result.kind.upper()} contains replacement-character mojibake.")


def _require_metric(result: ArtifactQaResult, key: str, minimum: int, message: str) -> None:
    value = int(result.metrics.get(key) or 0)
    if value < minimum:
        result.issues.append(f"{message} Observed {value}, expected at least {minimum}.")


def _overall_status(results: list[ArtifactQaResult]) -> str:
    if not results:
        return "failed"
    if any(item.structural_status == "failed" for item in results):
        return "failed"
    if any(item.render_status == "failed" for item in results):
        return "failed"
    if any(item.render_status == "blocked" for item in results):
        return "blocked"
    return "passed"


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="QA premium Smart Report DOCX/PPTX artifacts.")
    parser.add_argument("--docx", type=Path, help="Path to premium DOCX report.")
    parser.add_argument("--pptx", type=Path, help="Path to premium PPTX deck.")
    parser.add_argument("--out-dir", type=Path, default=Path("tmp/premium_artifact_qa"))
    parser.add_argument("--json", type=Path, help="Optional JSON output path.")
    parser.add_argument("--no-render", action="store_true", help="Skip LibreOffice/Poppler rendering.")
    parser.add_argument("--strict", action="store_true", help="Return non-zero for blocked visual render QA.")
    args = parser.parse_args(argv)

    report = run_qa(
        docx_path=args.docx,
        pptx_path=args.pptx,
        out_dir=args.out_dir,
        render=not args.no_render,
    )
    payload = json.dumps(report, ensure_ascii=False, indent=2)
    print(payload)
    if args.json:
        args.json.parent.mkdir(parents=True, exist_ok=True)
        args.json.write_text(payload + "\n", encoding="utf-8")
    if report["status"] == "failed":
        return 1
    if args.strict and report["status"] == "blocked":
        return 2
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
