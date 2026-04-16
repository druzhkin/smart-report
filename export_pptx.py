"""PowerPoint export: 10-15 minimalist slides. Reuses infographics from infographics.py.

Design: navy + white + accents, large key numbers, 1 idea per slide.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Cm, Inches, Pt, Emu

from infographics import render_all
from models import Block, BlockHeader, Connection, Report

NAVY = RGBColor(0x1B, 0x3A, 0x5C)
NAVY_LIGHT = RGBColor(0x2E, 0x59, 0x84)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x55, 0x55, 0x55)
GREY_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)
RED = RGBColor(0xC0, 0x39, 0x2B)
YELLOW = RGBColor(0xF1, 0xC4, 0x0F)
GREEN = RGBColor(0x27, 0xAE, 0x60)
BLUE = RGBColor(0x2E, 0x86, 0xC1)

PRIORITY_COLOR = {"high": RED, "medium": YELLOW, "low": GREEN}
PRIORITY_ORDER = {"high": 0, "medium": 1, "low": 2}
PRIORITY_LABEL = {"high": "🔴 Высокий", "medium": "🟡 Средний", "low": "🟢 Низкий"}

FONT = "Arial"

# Slide size: 13.333 × 7.5 inches (16:9 widescreen, default)


def _add_bg(slide, color: RGBColor = WHITE) -> None:
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0,
        slide.part.package.presentation_part.presentation.slide_width,
        slide.part.package.presentation_part.presentation.slide_height,
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()


def _add_textbox(
    slide, left, top, width, height, text: str,
    font_size: int = 18, bold: bool = False, color: RGBColor = NAVY,
    align=PP_ALIGN.LEFT, font: str = FONT, italic: bool = False,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.1); tf.margin_right = Cm(0.1)
    tf.margin_top = Cm(0.05); tf.margin_bottom = Cm(0.05)
    # first para
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def _add_accent_bar(slide, left, top, height, color: RGBColor = NAVY, width=Cm(0.2)) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def _sorted_blocks(report: Report) -> list[Block]:
    headers = {h.cell: h for h in report.block_headers}

    def _key(b: Block):
        h = headers.get(b.cell)
        prio = PRIORITY_ORDER.get(h.priority if h else "", 3)
        score = -(h.score_novelty + h.score_concreteness + h.score_applicability) if h else 0
        return (prio, score, b.cell)

    return sorted(report.blocks, key=_key)


# ---------- slide builders ----------

def _slide_title(prs, report: Report) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    sw, sh = prs.slide_width, prs.slide_height

    # navy left panel
    panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, sw // 3, sh)
    panel.fill.solid(); panel.fill.fore_color.rgb = NAVY
    panel.line.fill.background()

    # small label
    _add_textbox(
        slide, Cm(1), Cm(1), sw // 3 - Cm(2), Cm(1),
        "АНАЛИТИЧЕСКИЙ ОТЧЁТ", font_size=12, bold=True, color=WHITE,
    )

    # accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(1), Cm(2.2), Cm(3), Cm(0.08))
    line.fill.solid(); line.fill.fore_color.rgb = WHITE
    line.line.fill.background()

    # goal on right, big
    _add_textbox(
        slide, sw // 3 + Cm(1.5), Cm(3), sw - sw // 3 - Cm(3), Cm(6),
        report.goal, font_size=28, bold=True, color=NAVY,
    )

    # stats
    stats = (
        f"{len(report.matrix.domains)} доменов  ·  "
        f"{len(report.blocks)} блоков  ·  "
        f"{len(report.connections)} связей"
    )
    _add_textbox(
        slide, sw // 3 + Cm(1.5), sh - Cm(2.5), sw - sw // 3 - Cm(3), Cm(1),
        stats, font_size=14, color=GREY,
    )


def _slide_goal_context(prs, report: Report) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw = prs.slide_width
    _add_accent_bar(slide, Cm(1), Cm(1), Cm(1.5), NAVY)
    _add_textbox(slide, Cm(1.5), Cm(1), Cm(6), Cm(1), "ЦЕЛЬ И КОНТЕКСТ",
                 font_size=14, bold=True, color=NAVY)
    _add_textbox(slide, Cm(1), Cm(3), sw - Cm(2), Cm(3),
                 report.goal, font_size=24, bold=True, color=NAVY)
    es = report.exec_summary
    if es and es.goal_restate:
        _add_textbox(slide, Cm(1), Cm(8), sw - Cm(2), Cm(5),
                     es.goal_restate, font_size=14, color=GREY)


def _slide_image(prs, title: str, img_path: Path, subtitle: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw, sh = prs.slide_width, prs.slide_height
    _add_accent_bar(slide, Cm(1), Cm(1), Cm(1.5), NAVY)
    _add_textbox(slide, Cm(1.5), Cm(1), sw - Cm(3), Cm(1), title.upper(),
                 font_size=14, bold=True, color=NAVY)
    if subtitle:
        _add_textbox(slide, Cm(1), Cm(2), sw - Cm(2), Cm(1), subtitle,
                     font_size=12, color=GREY, italic=True)
    if img_path and Path(img_path).exists():
        # center image
        max_w = sw - Cm(3)
        max_h = sh - Cm(5)
        try:
            pic = slide.shapes.add_picture(str(img_path), Cm(1.5), Cm(3.5),
                                           width=max_w)
            # scale down if too tall
            if pic.height > max_h:
                ratio = max_h / pic.height
                pic.width = int(pic.width * ratio)
                pic.height = int(pic.height * ratio)
                pic.left = int((sw - pic.width) / 2)
                pic.top = int((sh - pic.height) / 2 + Cm(1))
        except Exception:
            pass


def _slide_finding(prs, finding_idx: int, headline, block: Block | None,
                   header: BlockHeader | None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw = prs.slide_width
    _add_accent_bar(slide, Cm(1), Cm(1), Cm(1.5), NAVY)
    _add_textbox(slide, Cm(1.5), Cm(1), Cm(10), Cm(1),
                 f"НАХОДКА #{finding_idx}",
                 font_size=14, bold=True, color=NAVY)

    # big headline
    _add_textbox(slide, Cm(1), Cm(3), sw - Cm(2), Cm(4),
                 headline.headline, font_size=22, bold=True, color=NAVY)

    # big number if available
    if header and header.strongest_number:
        _add_textbox(slide, Cm(1), Cm(9), sw - Cm(2), Cm(2.5),
                     header.strongest_number, font_size=40, bold=True, color=BLUE)
    # cell badge
    _add_textbox(slide, Cm(1), Cm(13), sw - Cm(2), Cm(1),
                 f"[{headline.block_cell}]", font_size=12,
                 color=GREY, italic=True)


def _slide_connection_detail(prs, conn_idx: int, conn: Connection) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw = prs.slide_width

    nature_icons = {
        "paradox": ("⚡ ПАРАДОКС", RED),
        "causal_chain": ("🔗 ПРИЧИННАЯ ЦЕПОЧКА", BLUE),
        "unexpected_confirmation": ("✓ ПОДТВЕРЖДЕНИЕ", GREEN),
        "shared_variable": ("◇ ОБЩАЯ ПЕРЕМЕННАЯ", NAVY_LIGHT),
    }
    label, color = nature_icons.get(conn.nature, (conn.nature.upper(), NAVY))

    _add_accent_bar(slide, Cm(1), Cm(1), Cm(1.5), color)
    _add_textbox(slide, Cm(1.5), Cm(1), sw - Cm(3), Cm(1),
                 f"СВЯЗЬ #{conn_idx} · {label}",
                 font_size=14, bold=True, color=color)

    # domains
    doms = " ↔ ".join(conn.domains)
    _add_textbox(slide, Cm(1), Cm(3), sw - Cm(2), Cm(1.5),
                 doms, font_size=20, bold=True, color=NAVY)

    # description
    _add_textbox(slide, Cm(1), Cm(5), sw - Cm(2), Cm(5),
                 conn.description, font_size=14, color=NAVY)

    # novelty (highlighted)
    if conn.novelty:
        _add_textbox(slide, Cm(1), Cm(11), sw - Cm(2), Cm(2.5),
                     f"Что нового: {conn.novelty}",
                     font_size=12, color=GREY, italic=True)
    # strength
    _add_textbox(slide, Cm(1), Cm(14), Cm(8), Cm(1),
                 f"Сила: {conn.strength}  ·  Общая переменная: {conn.shared_entity}",
                 font_size=10, color=GREY)


def _slide_gaps_next_steps(prs, report: Report) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw = prs.slide_width
    _add_accent_bar(slide, Cm(1), Cm(1), Cm(1.5), RED)
    _add_textbox(slide, Cm(1.5), Cm(1), sw - Cm(3), Cm(1),
                 "КРИТИЧЕСКИЕ ПРОБЕЛЫ И СЛЕДУЮЩИЕ ШАГИ",
                 font_size=14, bold=True, color=NAVY)
    es = report.exec_summary
    key_gaps = (es.key_gaps if es else [])[:5]
    if not key_gaps:
        # fall back to first gaps across blocks
        for b in report.blocks:
            for g in b.gaps[:2]:
                key_gaps.append(f"[{b.cell}] {g}")
                if len(key_gaps) >= 5:
                    break
            if len(key_gaps) >= 5:
                break

    y = Cm(3)
    for gap in key_gaps:
        bullet = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(1), y + Cm(0.2), Cm(0.4), Cm(0.4))
        bullet.fill.solid(); bullet.fill.fore_color.rgb = RED
        bullet.line.fill.background()
        _add_textbox(slide, Cm(2), y, sw - Cm(3), Cm(1.5),
                     gap, font_size=14, color=NAVY)
        y += Cm(1.8)


def _slide_analogies(prs, report: Report) -> None:
    """One slide with analogy cards for top-3 blocks that have analogies."""
    blocks_with = [b for b in _sorted_blocks(report) if getattr(b, "analogies", None)][:3]
    if not blocks_with:
        return
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw, sh = prs.slide_width, prs.slide_height
    _add_accent_bar(slide, Cm(1), Cm(1), Cm(1.5), NAVY)
    _add_textbox(slide, Cm(1.5), Cm(1), sw - Cm(3), Cm(1),
                 "АНАЛОГИИ К TOP-3 ВЫВОДАМ", font_size=12, bold=True, color=NAVY)
    card_w = (sw - Cm(2) - Cm(0.4) * (len(blocks_with) - 1)) // len(blocks_with)
    card_h = sh - Cm(4)
    for idx, b in enumerate(blocks_with):
        a = b.analogies[0]
        x = Cm(1) + idx * (card_w + Cm(0.4))
        y = Cm(2.5)
        # card background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, card_h)
        bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)
        bg.line.color.rgb = NAVY
        # location heading
        loc = getattr(a, "location", "") or b.cell
        _add_textbox(slide, x + Cm(0.3), y + Cm(0.2), card_w - Cm(0.6), Cm(0.8),
                     loc, font_size=11, bold=True, color=NAVY)
        # matched bullets
        matched = (getattr(a, "matched", None) or [])[:2]
        differed = (getattr(a, "differed", None) or [])[:2]
        content_lines: list[str] = []
        if matched:
            content_lines.append("+ " + matched[0])
            if len(matched) > 1:
                content_lines.append("+ " + matched[1])
        if differed:
            content_lines.append("~ " + differed[0])
            if len(differed) > 1:
                content_lines.append("~ " + differed[1])
        if content_lines:
            _add_textbox(slide, x + Cm(0.3), y + Cm(1.2), card_w - Cm(0.6), Cm(2.5),
                         "\n".join(content_lines), font_size=9, color=GREY)
        # lesson
        lesson_text = (getattr(a, "lesson", "") or "")
        _add_textbox(slide, x + Cm(0.3), y + card_h - Cm(1.6), card_w - Cm(0.6), Cm(1.4),
                     lesson_text, font_size=9, bold=True, color=NAVY)


def _slide_quadrant_crunch(prs, report: Report) -> None:
    """1 slide with top critical assumption inversions."""
    inv_blocks = getattr(report, "assumption_inversions", None) or []
    if not inv_blocks:
        return

    critical = [
        (bi.block_cell, inv)
        for bi in inv_blocks
        for inv in bi.inversions
        if inv.dependency == "critical"
    ][:4]

    unfalsifiable = [bi.block_cell for bi in inv_blocks if bi.unfalsifiable_flag]

    if not critical and not unfalsifiable:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw, sh = prs.slide_width, prs.slide_height
    _add_accent_bar(slide, Cm(1), Cm(1), Cm(1.5), NAVY)

    if not critical and unfalsifiable:
        title = "Внимание: нефальсифицируемые выводы"
        _add_textbox(slide, Cm(1.5), Cm(1), sw - Cm(3), Cm(1),
                     title, font_size=12, bold=True, color=NAVY)
        body = "Следующие блоки не имеют критически несущих допущений:\n" + "\n".join(unfalsifiable)
        _add_textbox(slide, Cm(1.5), Cm(2.5), sw - Cm(3), sh - Cm(4),
                     body, font_size=11, color=GREY)
        return

    _add_textbox(slide, Cm(1.5), Cm(1), sw - Cm(3), Cm(1),
                 "ПРОВЕРКА ДОПУЩЕНИЙ", font_size=12, bold=True, color=NAVY)

    col_w = (sw - Cm(2)) // 2
    left_x = Cm(1)
    right_x = Cm(1) + col_w + Cm(0.2)
    y_header = Cm(2.2)

    _add_textbox(slide, left_x, y_header, col_w - Cm(0.2), Cm(0.6),
                 "ДОПУЩЕНИЕ", font_size=8, bold=True, color=GREY)
    _add_textbox(slide, right_x, y_header, col_w - Cm(0.2), Cm(0.6),
                 "ЕСЛИ ЛОЖНО → ПОСЛЕДСТВИЕ", font_size=8, bold=True, color=RED)

    row_h = (sh - Cm(4)) // max(len(critical), 1)
    for i, (cell, inv) in enumerate(critical):
        y = Cm(3) + i * row_h
        left_text = f"{cell}\n{inv.assumption}"
        right_text = f"{inv.inversion}\n{inv.consequence}"
        _add_textbox(slide, left_x, y, col_w - Cm(0.2), row_h - Cm(0.2),
                     left_text, font_size=9, color=NAVY)
        _add_textbox(slide, right_x, y, col_w - Cm(0.2), row_h - Cm(0.2),
                     right_text, font_size=9, color=RED)


def _slide_scenarios(prs, report: Report) -> None:
    """1 slide for the Cone of Plausibility (predictive questions only)."""
    cone = getattr(report, "scenario_cone", None)
    if not cone:
        return
    scenarios = getattr(cone, "scenarios", []) or []
    if not scenarios:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw, sh = prs.slide_width, prs.slide_height
    horizon = getattr(cone, "question_horizon", "12-24 месяцев") or "12-24 месяцев"

    _add_accent_bar(slide, Cm(1), Cm(1), Cm(1.5), NAVY)
    _add_textbox(slide, Cm(1.5), Cm(1), sw - Cm(3), Cm(1),
                 f"ТРИ СЦЕНАРИЯ · горизонт {horizon}",
                 font_size=14, bold=True, color=NAVY)

    # verdict strip
    verdict = getattr(cone, "conditional_verdict", "") or ""
    if verdict:
        _add_textbox(slide, Cm(1), Cm(2.2), sw - Cm(2), Cm(1.0),
                     verdict, font_size=10, color=GREY, italic=True)

    # 3 columns
    n_cols = min(len(scenarios), 3)
    col_colors = [BLUE, GREEN, RED]  # base, optimistic, pessimistic defaults
    gap = Cm(0.3)
    col_w = (sw - Cm(2) - gap * (n_cols - 1)) // n_cols
    card_top = Cm(3.5)
    card_h = sh - card_top - Cm(1.5)

    for idx, s in enumerate(scenarios[:n_cols]):
        name = getattr(s, "name", "") or ""
        prob = getattr(s, "probability", "") or ""
        desc = getattr(s, "description", "") or ""
        implications = (getattr(s, "implications", []) or [])[:2]
        indicators = (getattr(s, "indicators", []) or [])[:2]

        name_low = name.lower()
        if "optim" in name_low or "оптим" in name_low:
            accent = GREEN
        elif "pessim" in name_low or "пессим" in name_low:
            accent = RED
        else:
            accent = BLUE

        x = Cm(1) + idx * (col_w + gap)
        # card bg
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, card_top, col_w, card_h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = GREY_LIGHT
        bg.line.color.rgb = accent
        bg.line.width = Pt(1.5)

        # probability (big)
        _add_textbox(slide, x + Cm(0.3), card_top + Cm(0.2), col_w - Cm(0.6), Cm(1),
                     prob, font_size=22, bold=True, color=accent)
        # name
        _add_textbox(slide, x + Cm(0.3), card_top + Cm(1.1), col_w - Cm(0.6), Cm(0.6),
                     name, font_size=10, bold=True, color=NAVY)
        # description (truncated)
        _add_textbox(slide, x + Cm(0.3), card_top + Cm(1.8), col_w - Cm(0.6), Cm(2.0),
                     (desc[:180] + "…" if len(desc) > 180 else desc), font_size=9, color=GREY)
        # implications
        if implications:
            impl_text = "\n".join(f"› {i}" for i in implications)
            _add_textbox(slide, x + Cm(0.3), card_top + Cm(3.9), col_w - Cm(0.6), Cm(1.6),
                         impl_text, font_size=8, color=NAVY)
        # indicators
        if indicators:
            ind_text = "\n".join(f"◎ {i}" for i in indicators)
            _add_textbox(slide, x + Cm(0.3), card_top + Cm(5.6), col_w - Cm(0.6), Cm(1.4),
                         ind_text, font_size=7, color=GREY)

    # wild card footer
    wc = getattr(cone, "wild_card", None)
    if wc:
        wc_text = f"Wild Card ({getattr(wc, 'probability', '')}) — {(getattr(wc, 'description', '') or '')[:100]}"
        _add_textbox(slide, Cm(1), sh - Cm(1.2), sw - Cm(2), Cm(1),
                     wc_text, font_size=8, color=RGBColor(0xD9, 0x77, 0x06), italic=True)


def _slide_conclusion(prs, report: Report) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw, sh = prs.slide_width, prs.slide_height

    # navy bg
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, sw, sh)
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    _add_textbox(slide, Cm(1), Cm(3), sw - Cm(2), Cm(2),
                 "ЗАКЛЮЧЕНИЕ", font_size=14, bold=True, color=WHITE)

    # stats line
    stats = (
        f"{len(report.matrix.domains)} доменов · {len(report.blocks)} блоков · "
        f"{len(report.connections)} связей"
    )
    _add_textbox(slide, Cm(1), Cm(6), sw - Cm(2), Cm(3),
                 stats, font_size=26, bold=True, color=WHITE)

    # top finding echo
    es = report.exec_summary
    if es and es.top_findings:
        _add_textbox(slide, Cm(1), Cm(11), sw - Cm(2), Cm(3),
                     es.top_findings[0].headline,
                     font_size=18, color=WHITE, italic=True)


# ---------- main ----------

def export_pptx(report: Report, path: Path, images: dict[str, Path] | None = None) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    if images is None:
        images = render_all(report)

    # 1. Title
    _slide_title(prs, report)
    # 2. Goal and context
    _slide_goal_context(prs, report)
    # 3. Matrix map
    _slide_image(prs, "Матрица доменов", images.get("matrix"),
                 "Цвет = приоритет домена (🔴/🟡/🟢)")
    # 4. Key Metrics Dashboard
    _slide_image(prs, "Ключевые метрики", images.get("metrics"),
                 "Самые сильные цифры отчёта")
    # 5. Top findings — up to 5 slides
    es = report.exec_summary
    headers = {h.cell: h for h in report.block_headers}
    blocks_by_cell = {b.cell: b for b in report.blocks}
    if es and es.top_findings:
        for i, tf in enumerate(es.top_findings[:5], start=1):
            _slide_finding(prs, i, tf,
                           blocks_by_cell.get(tf.block_cell),
                           headers.get(tf.block_cell))
    # 6. Connections graph
    _slide_image(prs, "Граф кросс-доменных связей", images.get("graph"),
                 "Цвет ребра = тип связи, толщина = сила")
    # 7. Top-3 connection details
    top_conns = report.connections[:3]
    for i, c in enumerate(top_conns, start=1):
        _slide_connection_detail(prs, i, c)
    # 8. Heatmap
    _slide_image(prs, "Тепловая карта приоритетов", images.get("heatmap"),
                 "Где золото, где пусто")
    # 9. Analogies
    _slide_analogies(prs, report)
    # 10. Quadrant Crunch — assumption inversions
    _slide_quadrant_crunch(prs, report)
    # 11. Scenarios (predictive questions only)
    _slide_scenarios(prs, report)
    # 11. Gaps & next steps
    _slide_gaps_next_steps(prs, report)
    # 12. Conclusion
    _slide_conclusion(prs, report)

    path = Path(path)
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))
    return path
