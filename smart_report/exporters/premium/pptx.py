"""Premium editable PPTX renderer.

This is a separate deck artifact for the premium package. It uses the same
renderer-neutral PremiumReportDocument as the DOCX report, but optimizes the
output for presentation: answer, evidence, readiness, scenarios/risks, and
section-level implications. It is intentionally domain-neutral.
"""

from __future__ import annotations

from pathlib import Path

from .models import PremiumPreparedBlock, PremiumPreparedSection, PremiumReportDocument

NAVY = "152238"
INK = "1F2933"
MUTED = "667085"
GOLD = "B08D57"
PAPER = "F7F4EE"
LINE = "D9DEE7"
RED = "B42318"
GREEN = "027A48"


def render_premium_pptx(document: PremiumReportDocument, path: Path) -> Path:
    """Render an editable premium PPTX deck."""

    from pptx import Presentation
    from pptx.util import Inches, Pt

    path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = document.title
    prs.core_properties.subject = document.subtitle
    prs.core_properties.author = "Smart Report"

    blank = prs.slide_layouts[6]
    _cover(prs, blank, document)
    _executive_answer(prs, blank, document)
    _readiness(prs, blank, document)
    _evidence(prs, blank, document)

    for idx, section in enumerate(document.sections[:6], start=1):
        _section_slide(prs, blank, idx, section)

    _deck_close(prs, blank, document)
    prs.save(str(path))
    return path


def _cover(prs, layout, document: PremiumReportDocument) -> None:
    slide = prs.slides.add_slide(layout)
    _brand_bar(slide)
    _textbox(slide, 0.65, 0.55, 4.0, 0.3, "SMART REPORT", 8, GOLD, bold=True)
    _textbox(slide, 0.65, 1.55, 10.4, 1.15, document.title, 28, NAVY, bold=True)
    _textbox(slide, 0.68, 2.82, 8.8, 0.55, document.subtitle, 13, MUTED)
    _metric_card(slide, 0.7, 4.65, "Report", document.plan.report_type.replace("_", " ").title())
    _metric_card(slide, 3.2, 4.65, "Audience", document.plan.audience.replace("_", " ").title())
    _metric_card(slide, 5.7, 4.65, "Sources", str(document.source_count))
    _metric_card(slide, 8.2, 4.65, "Numeric facts", str(document.numeric_fact_count))
    _footer(slide, 1)


def _executive_answer(prs, layout, document: PremiumReportDocument) -> None:
    slide = prs.slides.add_slide(layout)
    _title(slide, "Executive Answer")
    answer = _first_block_body(document.sections, "executive_summary", "Answer")
    _textbox(slide, 0.7, 1.25, 7.2, 1.45, answer or document.plan.decision_context, 18, INK, bold=True)
    _table(
        slide,
        0.7,
        3.15,
        5.9,
        2.45,
        ["Decision lens", "Implication"],
        [
            ["Evidence", f"{document.source_count} sources / {document.numeric_fact_count} facts"],
            ["Quality bar", document.plan.quality_bar],
            ["Deliverables", _deliverables(document)],
        ],
    )
    _callout(slide, 7.1, 3.1, 4.9, 2.5, document.plan.decision_context)
    _footer(slide, 2)


def _readiness(prs, layout, document: PremiumReportDocument) -> None:
    slide = prs.slides.add_slide(layout)
    _title(slide, "Paid-Delivery Readiness")
    readiness = document.premium_readiness or {}
    ready = bool(readiness.get("ready"))
    score = readiness.get("score", "?")
    status = "READY" if ready else "NOT READY"
    color = GREEN if ready else RED
    _metric_card(slide, 0.7, 1.25, "Status", status, value_color=color, width=2.6)
    _metric_card(slide, 3.55, 1.25, "Score", f"{score}/100", value_color=color, width=2.6)
    issues = [issue for issue in readiness.get("issues", []) if isinstance(issue, dict)]
    _metric_card(slide, 6.4, 1.25, "Open issues", str(len(issues)), value_color=color, width=2.6)
    rows = [
        [
            str(issue.get("severity", "")),
            str(issue.get("message", ""))[:95],
            str(issue.get("recommendation", ""))[:95],
        ]
        for issue in issues[:5]
    ] or [["-", "No premium readiness issues in the exported gate.", "-"]]
    _table(slide, 0.7, 3.0, 11.8, 3.05, ["Severity", "Issue", "Recommended fix"], rows)
    _footer(slide, 3)


def _evidence(prs, layout, document: PremiumReportDocument) -> None:
    slide = prs.slides.add_slide(layout)
    _title(slide, "Evidence Base")
    rows = []
    for section in document.sections:
        for block in section.blocks:
            if block.kind in {"evidence_table", "source_quality_table", "kpi_grid"}:
                rows.append([block.title, block.kind.replace("_", " "), str(len(block.rows))])
    _table(slide, 0.7, 1.35, 11.6, 4.85, ["Block", "Type", "Rows"], rows[:8])
    _footer(slide, 4)


def _section_slide(prs, layout, number: int, section: PremiumPreparedSection) -> None:
    slide = prs.slides.add_slide(layout)
    _title(slide, section.title)
    _textbox(slide, 0.72, 1.05, 10.9, 0.5, section.purpose, 12, MUTED)
    rows = []
    for block in section.blocks[:4]:
        rows.append([block.title, block.kind.replace("_", " "), _block_signal(block)])
    _table(slide, 0.7, 1.95, 11.6, 3.8, ["Analytical block", "Format", "Signal"], rows)
    _footer(slide, number + 4)


def _deck_close(prs, layout, document: PremiumReportDocument) -> None:
    slide = prs.slides.add_slide(layout)
    _title(slide, "Next Decisions")
    rows = [
        ["1", "Close evidence blockers", "Run targeted follow-up for critical gaps and disputed figures."],
        ["2", "Lock client narrative", "Convert the full report into a short board-level storyline."],
        ["3", "Package delivery", _deliverables(document)],
    ]
    _table(slide, 0.7, 1.45, 11.6, 3.3, ["#", "Move", "Why it matters"], rows)
    _callout(
        slide,
        0.7,
        5.3,
        11.6,
        0.75,
        "The deck is an executive presentation. The DOCX report remains the full evidence-backed deliverable.",
    )
    _footer(slide, len(document.sections[:6]) + 5)


def _brand_bar(slide) -> None:
    from pptx.util import Inches

    vertical = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.12), Inches(7.5))
    vertical.fill.solid()
    vertical.fill.fore_color.rgb = _rgb(GOLD)
    horizontal = slide.shapes.add_shape(1, Inches(0.12), Inches(0), Inches(13.213), Inches(0.12))
    horizontal.fill.solid()
    horizontal.fill.fore_color.rgb = _rgb(NAVY)


def _title(slide, text: str) -> None:
    _brand_bar(slide)
    _textbox(slide, 0.68, 0.38, 10.5, 0.55, text, 22, NAVY, bold=True)


def _footer(slide, page: int) -> None:
    _textbox(slide, 0.7, 6.92, 6.0, 0.25, "Smart Report | Premium deck", 7, MUTED)
    _textbox(slide, 11.8, 6.92, 0.8, 0.25, str(page), 7, MUTED)


def _metric_card(slide, x: float, y: float, label: str, value: str, *, value_color: str = NAVY, width: float = 2.25) -> None:
    from pptx.util import Inches

    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(width), Inches(0.92))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(PAPER)
    shape.line.color.rgb = _rgb(LINE)
    _textbox(slide, x + 0.12, y + 0.12, width - 0.24, 0.18, label, 7, MUTED, bold=True)
    _textbox(slide, x + 0.12, y + 0.36, width - 0.24, 0.38, value, 13, value_color, bold=True)


def _callout(slide, x: float, y: float, w: float, h: float, text: str) -> None:
    from pptx.util import Inches

    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(NAVY)
    shape.line.color.rgb = _rgb(NAVY)
    _textbox(slide, x + 0.18, y + 0.18, w - 0.36, h - 0.25, text, 12, "FFFFFF", bold=True)


def _table(slide, x: float, y: float, w: float, h: float, headers: list[str], rows: list[list[str]]) -> None:
    from pptx.util import Inches, Pt

    rows = rows or [["-", "-", "-"]]
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(NAVY)
        cell.text = header
        _cell_text(cell, color="FFFFFF", bold=True, size=8)
    for row_idx, row in enumerate(rows, start=1):
        for col_idx, value in enumerate(row[: len(headers)]):
            cell = table.cell(row_idx, col_idx)
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb("FAFBFC")
            cell.text = str(value or "")
            _cell_text(cell, color=INK, size=7.5)
    for row in table.rows:
        for cell in row.cells:
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)


def _textbox(slide, x: float, y: float, w: float, h: float, text: str, size: float, color: str, *, bold: bool = False) -> None:
    from pptx.util import Inches, Pt

    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = str(text or "")
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.color.rgb = _rgb(color)
    run.font.bold = bold


def _cell_text(cell, *, color: str, size: float = 8, bold: bool = False) -> None:
    from pptx.util import Pt

    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = "Aptos"
            run.font.size = Pt(size)
            run.font.color.rgb = _rgb(color)
            run.font.bold = bold


def _rgb(hex_color: str):
    from pptx.dml.color import RGBColor

    value = hex_color.strip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def _first_block_body(sections: list[PremiumPreparedSection], section_id: str, block_title: str) -> str:
    for section in sections:
        if section.id != section_id:
            continue
        for block in section.blocks:
            if block.title == block_title:
                return block.body
    return ""


def _block_signal(block: PremiumPreparedBlock) -> str:
    if block.rows:
        return f"{len(block.rows)} row(s)"
    if block.body:
        return block.body[:120]
    if block.notes:
        return block.notes[0][:120]
    return "Prepared for narrative synthesis"


def _deliverables(document: PremiumReportDocument) -> str:
    deliverables = document.plan.deliverables
    names = []
    if deliverables.require_docx:
        names.append("DOCX report")
    if deliverables.require_pptx:
        names.append("PPTX deck")
    if deliverables.require_pdf:
        names.append("PDF")
    if deliverables.require_data_pack:
        names.append("data pack")
    if deliverables.require_qa_audit:
        names.append("QA audit")
    return ", ".join(names)
