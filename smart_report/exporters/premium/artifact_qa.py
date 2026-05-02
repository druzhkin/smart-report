"""QA checks for premium Smart Report DOCX, PDF, and PPTX artifacts.

The checker is intentionally additive: it does not change exported files. It
performs structural inspection everywhere and performs visual rendering only
when LibreOffice (`soffice`) and Poppler (`pdftoppm`) are available locally.
"""

from __future__ import annotations

import argparse
import html
import json
import re
import shutil
import struct
import subprocess
import sys
import tempfile
import zlib
from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Any

INTERNAL_MARKERS = (
    "[STRONG]",
    "[MODERATE]",
    "[WEAK]",
    "[SPECULATIVE]",
    "[REF:",
    "main_synthesis",
    "consensus_section",
    "gaps_filled_section",
    "source_reports",
    "followup_reports",
)

DOCX_NARRATIVE_MIN_CHARS = 1800
DOCX_MAX_PARAGRAPH_CHARS = 650
PDF_THIN_PAGE_MIN_CHARS = 180
PLACEHOLDER_CONTENT_PATTERNS = (
    "example.com",
    "internal QA fixture",
    "stub source",
    "demo evidence",
    "placeholder",
    "Forecast a market with scenario and risk recommendations",
)
SOURCE_REFERENCE_RE = re.compile(
    r"https?://|www\.|(?:^|[^A-Za-z])doi[:/ ]|\b[\w.-]+\.(?:com|ru|org|net|io|gov|edu|рф)\b",
    re.IGNORECASE,
)

VISUAL_REVIEW_RUBRIC = (
    "No text overflow, clipping, or overlapping elements.",
    "Tables are readable without broken rows, orphan headers, or cramped cells.",
    "Headings create a clear executive-to-detail hierarchy.",
    "Page and slide breaks preserve complete ideas and do not strand captions.",
    "Charts, scorecards, and badges are visually aligned and easy to scan.",
    "The package looks like a finished paid report, not a raw model export.",
)

RENDER_BLANK_INK_RATIO = 0.002
RENDER_LOW_INK_RATIO = 0.01


@dataclass
class ArtifactQaResult:
    path: str
    kind: str
    exists: bool
    structural_status: str = "not_run"
    render_status: str = "not_run"
    missing_tools: list[str] = field(default_factory=list)
    metrics: dict[str, Any] = field(default_factory=dict)
    issues: list[str] = field(default_factory=list)
    rendered_files: list[str] = field(default_factory=list)

    @property
    def passed(self) -> bool:
        return self.exists and self.structural_status == "passed" and not self.issues


def run_qa(
    *,
    docx_path: Path | None = None,
    pdf_path: Path | None = None,
    pptx_path: Path | None = None,
    out_dir: Path,
    render: bool = True,
) -> dict[str, Any]:
    out_dir.mkdir(parents=True, exist_ok=True)
    results: list[ArtifactQaResult] = []
    if docx_path:
        results.append(_inspect_docx(docx_path))
    if pdf_path:
        results.append(_inspect_pdf(pdf_path))
    if pptx_path:
        results.append(_inspect_pptx(pptx_path))

    if render:
        soffice = _find_tool(
            "soffice",
            [
                Path("C:/Program Files/LibreOffice/program/soffice.exe"),
                Path("C:/Program Files (x86)/LibreOffice/program/soffice.exe"),
            ],
        )
        pdftoppm = _find_tool(
            "pdftoppm",
            list(Path.home().glob(
                "AppData/Local/Microsoft/WinGet/Packages/oschwartz10612.Poppler_*/poppler-*/Library/bin/pdftoppm.exe"
            )),
        )
        for result in results:
            _render_artifact(result, out_dir, soffice=soffice, pdftoppm=pdftoppm)
    else:
        for result in results:
            result.render_status = "skipped"

    report = {
        "status": _overall_status(results),
        "summary": {
            "artifacts": len(results),
            "passed_structural": sum(1 for item in results if item.structural_status == "passed"),
            "rendered": sum(1 for item in results if item.render_status == "passed"),
            "blocked_render": sum(1 for item in results if item.render_status == "blocked"),
            "issues": sum(len(item.issues) for item in results),
        },
        "results": [asdict(item) for item in results],
        "manual_visual_review_rubric": list(VISUAL_REVIEW_RUBRIC),
    }
    render_index = _write_render_index(results, out_dir)
    if render_index:
        report["render_index"] = str(render_index)
    return report


def _inspect_docx(path: Path) -> ArtifactQaResult:
    result = ArtifactQaResult(path=str(path), kind="docx", exists=path.exists())
    if not result.exists:
        result.structural_status = "failed"
        result.issues.append("DOCX file does not exist.")
        return result

    try:
        from docx import Document

        doc = Document(path)
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        table_text = [
            cell.text.strip()
            for table in doc.tables
            for row in table.rows
            for cell in row.cells
            if cell.text.strip()
        ]
        text = "\n".join([*paragraphs, *table_text])
        headings = [
            p.text.strip()
            for p in doc.paragraphs
            if p.style and p.style.name.startswith("Heading") and p.text.strip()
        ]
        narrative_paragraphs = [
            p.text.strip()
            for p in doc.paragraphs
            if p.text.strip()
            and not (p.style and p.style.name.startswith(("Heading", "Title")))
        ]
        overlong_paragraphs = [
            {"index": idx, "chars": len(paragraph)}
            for idx, paragraph in enumerate(narrative_paragraphs, start=1)
            if len(paragraph) > DOCX_MAX_PARAGRAPH_CHARS
        ]
        source_reference_count = len(SOURCE_REFERENCE_RE.findall(text))
        result.metrics = {
            "paragraphs": len(paragraphs),
            "narrative_paragraphs": len(narrative_paragraphs),
            "tables": len(doc.tables),
            "sections": len(doc.sections),
            "headings": len(headings),
            "text_chars": len(text),
            "narrative_chars": sum(len(paragraph) for paragraph in narrative_paragraphs),
            "source_reference_count": source_reference_count,
            "overlong_paragraphs": overlong_paragraphs,
            "estimated_pages": _estimate_docx_pages(text, len(doc.tables)),
            "has_cover_brand": (
                "SMART REPORT | PREMIUM ANALYTICAL REPORT" in text
                or "SMART REPORT | ПРЕМИАЛЬНЫЙ АНАЛИТИЧЕСКИЙ ОТЧЁТ" in text
            ),
            "has_decision_dashboard": (
                "Client Decision Dashboard" in text
                or "Панель решения клиента" in text
                or "Резюме для решения" in text
            ),
            "has_scorecard": (
                "Executive Evidence Scorecard" in text
                or "Карта доказательной базы" in text
                or "Карта доказательств" in text
            ),
            "has_readiness_gate": (
                "Premium Readiness Gate" in text or "Гейт готовности к платной выдаче" in text
            ),
            "has_report_structure": (
                "Report Structure" in text
                or "How to Read" in text
                or "Структура отчёта" in text
                or "Как читать отчёт" in text
            ),
        }
        _add_common_content_issues(result, text)
        _require_metric(result, "paragraphs", 25, "DOCX has too few paragraphs for a long-form report.")
        _require_metric(
            result,
            "narrative_chars",
            DOCX_NARRATIVE_MIN_CHARS,
            "DOCX has too little narrative prose outside headings and tables.",
        )
        _require_metric(result, "tables", 4, "DOCX has too few visual tables.")
        _require_metric(result, "text_chars", 5000, "DOCX text volume is below structural QA minimum.")
        _require_metric(
            result,
            "source_reference_count",
            1,
            "DOCX has no traceable source references.",
        )
        if overlong_paragraphs:
            sample = ", ".join(
                f"#{item['index']}={item['chars']} chars" for item in overlong_paragraphs[:5]
            )
            result.issues.append(
                "DOCX contains overlong paragraph(s); split text with bullets, callouts, "
                f"or exhibits: {sample}."
            )
        for key, message in {
            "has_cover_brand": "DOCX cover brand marker is missing.",
            "has_decision_dashboard": "DOCX decision summary is missing.",
            "has_scorecard": "DOCX evidence map is missing.",
            "has_report_structure": "DOCX report structure section is missing.",
        }.items():
            if not result.metrics.get(key):
                result.issues.append(message)
        result.structural_status = "passed" if not result.issues else "failed"
    except Exception as exc:  # pragma: no cover - defensive CLI path
        result.structural_status = "failed"
        result.issues.append(f"DOCX inspection failed: {exc}")
    return result


def _estimate_docx_pages(text: str, table_count: int) -> int:
    """Approximate page count when render tools are unavailable.

    It is deliberately conservative and only a metric; real paid-delivery page
    QA still requires DOCX/PPTX rendering through LibreOffice/Poppler.
    """
    text_pages = len(text) / 2800
    table_pages = table_count * 0.15
    return max(1, round(text_pages + table_pages))


def _inspect_pdf(path: Path) -> ArtifactQaResult:
    result = ArtifactQaResult(path=str(path), kind="pdf", exists=path.exists())
    if not result.exists:
        result.structural_status = "failed"
        result.issues.append("PDF file does not exist.")
        return result

    try:
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        page_count = len(reader.pages)
        page_texts = [(page.extract_text() or "") for page in reader.pages]
        full_text = "\n".join(page_texts)
        sample_text = "\n".join(
            page_texts[: min(page_count, 6)]
        )
        page_char_counts = [len(text.strip()) for text in page_texts]
        placeholder_hits = sorted(
            {
                pattern
                for pattern in PLACEHOLDER_CONTENT_PATTERNS
                if pattern.lower() in full_text.lower()
            }
        )
        landscape_pages = []
        thin_pages = []
        sparse_pages = []
        pages_without_source_note = []
        max_text_only_streak = 0
        current_text_only_streak = 0
        for idx, page in enumerate(reader.pages, start=1):
            width = float(page.mediabox.width)
            height = float(page.mediabox.height)
            if width > height:
                landscape_pages.append(idx)
            # Skip the cover for density checks: a publication cover can be sparse.
            if idx > 1 and page_char_counts[idx - 1] < PDF_THIN_PAGE_MIN_CHARS:
                thin_pages.append(idx)
            page_text = page_texts[idx - 1]
            if idx > 1 and PDF_THIN_PAGE_MIN_CHARS <= page_char_counts[idx - 1] < 360:
                sparse_pages.append(idx)
            if idx > 1 and page_char_counts[idx - 1] > 250 and not _has_source_note(page_text):
                pages_without_source_note.append(idx)
            page_is_exhibit = _has_exhibit_marker(page_text)
            if idx > 1 and page_char_counts[idx - 1] > 900 and not page_is_exhibit:
                current_text_only_streak += 1
                max_text_only_streak = max(max_text_only_streak, current_text_only_streak)
            else:
                current_text_only_streak = 0
        result.metrics = {
            "pages": page_count,
            "text_chars_sample": len(sample_text),
            "page_text_chars_min": min(page_char_counts) if page_char_counts else 0,
            "page_text_chars_avg": round(sum(page_char_counts) / page_count) if page_count else 0,
            "landscape_pages": landscape_pages,
            "thin_pages_after_cover": thin_pages,
            "sparse_pages_after_cover": sparse_pages,
            "pages_without_source_note": pages_without_source_note,
            "max_dense_text_streak_after_cover": max_text_only_streak,
            "placeholder_content_hits": placeholder_hits,
            "has_cover_brand": "SMART REPORT" in sample_text,
            "has_publication_marker": "Publication-grade PDF" in sample_text,
            "has_exhibit_pages": "EXHIBIT" in sample_text,
            "has_source_notes": "Source:" in sample_text or "Источник:" in sample_text,
        }
        _add_common_content_issues(result, full_text)
        _require_metric(result, "pages", 20, "PDF has too few pages for a publication-grade report.")
        _require_metric(result, "text_chars_sample", 500, "PDF text extraction sample is too thin.")
        for key, message in {
            "has_cover_brand": "PDF cover brand marker is missing.",
            "has_publication_marker": "PDF publication marker is missing.",
            "has_exhibit_pages": "PDF exhibit pages are missing.",
            "has_source_notes": "PDF exhibit source notes are missing.",
        }.items():
            if not result.metrics.get(key):
                result.issues.append(message)
        if landscape_pages:
            result.issues.append(
                "PDF contains landscape page(s); publication report must be portrait: "
                + ", ".join(str(page) for page in landscape_pages[:10])
            )
        if placeholder_hits:
            result.issues.append(
                "PDF contains placeholder or demo content that cannot be delivered to a client: "
                + ", ".join(placeholder_hits[:10])
            )
        if thin_pages:
            result.issues.append(
                "PDF contains empty or underwritten page(s) after the cover: "
                + ", ".join(str(page) for page in thin_pages[:10])
            )
        if len(sparse_pages) >= 3:
            result.issues.append(
                "PDF has multiple sparse pages that may read like weak slides, not a report: "
                + ", ".join(str(page) for page in sparse_pages[:10])
            )
        if len(pages_without_source_note) >= max(3, page_count // 3):
            result.issues.append(
                "PDF has too many content pages without visible source notes: "
                + ", ".join(str(page) for page in pages_without_source_note[:10])
            )
        if max_text_only_streak >= 4:
            result.issues.append(
                "PDF has too many consecutive dense narrative pages without exhibit pacing; "
                f"observed streak {max_text_only_streak}, expected below 4."
            )
        result.structural_status = "passed" if not result.issues else "failed"
    except Exception as exc:  # pragma: no cover - defensive CLI path
        result.structural_status = "failed"
        result.issues.append(f"PDF inspection failed: {exc}")
    return result


def _has_source_note(text: str) -> bool:
    lowered = (text or "").lower()
    return any(marker in lowered for marker in ("source:", "sources:", "источник:", "источники:"))


def _has_exhibit_marker(text: str) -> bool:
    lowered = (text or "").lower()
    return any(marker in lowered for marker in ("exhibit", "рисунок", "figure", "appendix"))


def _find_tool(name: str, candidates: list[Path]) -> str | None:
    found = shutil.which(name)
    if found:
        return found
    for candidate in candidates:
        if candidate.exists():
            return str(candidate)
    return None


def _inspect_pptx(path: Path) -> ArtifactQaResult:
    result = ArtifactQaResult(path=str(path), kind="pptx", exists=path.exists())
    if not result.exists:
        result.structural_status = "failed"
        result.issues.append("PPTX file does not exist.")
        return result

    try:
        from pptx import Presentation

        deck = Presentation(str(path))
        slide_texts: list[str] = []
        shape_count = 0
        table_count = 0
        for slide in deck.slides:
            for shape in slide.shapes:
                shape_count += 1
                if getattr(shape, "has_table", False):
                    table_count += 1
                if hasattr(shape, "text") and shape.text:
                    slide_texts.append(shape.text.strip())
        text = "\n".join(part for part in slide_texts if part)
        result.metrics = {
            "slides": len(deck.slides),
            "shapes": shape_count,
            "tables": table_count,
            "text_chars": len(text),
            "has_executive_answer": "Executive Answer" in text or "Короткий ответ" in text,
            "has_client_position": (
                "Client Position" in text
                or "Позиция и ограничения" in text
            ),
            "has_readiness": (
                "Paid-Delivery Readiness" in text or "Готовность к платной выдаче" in text
            ),
            "has_evidence_base": "Evidence Base" in text or "Доказательная база" in text,
        }
        _add_common_content_issues(result, text)
        _require_metric(result, "slides", 10, "PPTX has too few slides for the premium deck.")
        _require_metric(result, "tables", 3, "PPTX has too few table-based analytical slides.")
        _require_metric(result, "text_chars", 1000, "PPTX text volume is below structural QA minimum.")
        for key, message in {
            "has_executive_answer": "PPTX executive answer slide is missing.",
            "has_client_position": "PPTX client position slide is missing.",
            "has_evidence_base": "PPTX evidence base slide is missing.",
        }.items():
            if not result.metrics.get(key):
                result.issues.append(message)
        result.structural_status = "passed" if not result.issues else "failed"
    except Exception as exc:  # pragma: no cover - defensive CLI path
        result.structural_status = "failed"
        result.issues.append(f"PPTX inspection failed: {exc}")
    return result


def _render_artifact(
    result: ArtifactQaResult,
    out_dir: Path,
    *,
    soffice: str | None,
    pdftoppm: str | None,
) -> None:
    if not result.exists:
        result.render_status = "failed"
        return
    missing = []
    if result.kind != "pdf" and not soffice:
        missing.append("soffice")
    if not pdftoppm:
        missing.append("pdftoppm")
    if missing:
        result.render_status = "blocked"
        result.missing_tools = missing
        result.issues.append(
            "Visual render QA was not completed because required tools are missing: "
            + ", ".join(missing)
            + "."
        )
        return

    artifact = Path(result.path)
    artifact_out = out_dir / f"{artifact.stem}_{result.kind}"
    artifact_out.mkdir(parents=True, exist_ok=True)
    if result.kind == "pdf":
        prefix = artifact_out / artifact.stem
        try:
            subprocess.run(
                [pdftoppm, "-png", str(artifact), str(prefix)],
                check=True,
                capture_output=True,
                text=True,
                timeout=120,
            )
            rendered = sorted(str(item) for item in artifact_out.glob(f"{artifact.stem}-*.png"))
            result.rendered_files = [str(artifact), *rendered]
            result.render_status = "passed" if rendered else "failed"
            if rendered:
                result.metrics["rendered_pages"] = len(rendered)
                result.metrics["rendered_total_bytes"] = _total_file_size(rendered)
                _add_render_density_metrics(result, rendered, page_label="page")
            else:
                result.issues.append("PDF rendering produced no PNG pages.")
        except subprocess.SubprocessError as exc:
            result.render_status = "failed"
            result.issues.append(f"Visual render command failed: {exc}")
        return

    profile = Path(tempfile.mkdtemp(prefix="smart-report-lo-"))
    try:
        subprocess.run(
            [
                soffice,
                f"-env:UserInstallation=file:///{profile.as_posix()}",
                "--headless",
                "--convert-to",
                "pdf",
                "--outdir",
                str(artifact_out),
                str(artifact),
            ],
            check=True,
            capture_output=True,
            text=True,
            timeout=120,
        )
        pdf_path = artifact_out / f"{artifact.stem}.pdf"
        if not pdf_path.exists():
            result.render_status = "failed"
            result.issues.append("LibreOffice conversion did not produce a PDF.")
            return
        prefix = artifact_out / artifact.stem
        subprocess.run(
            [pdftoppm, "-png", str(pdf_path), str(prefix)],
            check=True,
            capture_output=True,
            text=True,
            timeout=120,
        )
        rendered = sorted(str(item) for item in artifact_out.glob(f"{artifact.stem}-*.png"))
        result.rendered_files = [str(pdf_path), *rendered]
        result.render_status = "passed" if rendered else "failed"
        if rendered:
            result.metrics["rendered_total_bytes"] = _total_file_size(rendered)
            _add_render_density_metrics(
                result,
                rendered,
                page_label="page" if result.kind == "docx" else "slide",
            )
            if result.kind == "docx":
                result.metrics["rendered_pages"] = len(rendered)
            elif result.kind == "pptx":
                result.metrics["rendered_slides"] = len(rendered)
        if not rendered:
            result.issues.append("PDF rendering produced no PNG pages.")
    except subprocess.SubprocessError as exc:
        result.render_status = "failed"
        result.issues.append(f"Visual render command failed: {exc}")
    finally:
        shutil.rmtree(profile, ignore_errors=True)


def _total_file_size(paths: list[str]) -> int:
    total = 0
    for path in paths:
        try:
            total += Path(path).stat().st_size
        except OSError:
            continue
    return total


def _add_render_density_metrics(
    result: ArtifactQaResult,
    rendered_paths: list[str],
    *,
    page_label: str,
) -> None:
    densities: list[float] = []
    blank_pages: list[int] = []
    low_ink_pages: list[int] = []
    unreadable_pages: list[int] = []
    for idx, path in enumerate(rendered_paths, start=1):
        ratio = _png_ink_ratio(Path(path))
        if ratio is None:
            unreadable_pages.append(idx)
            continue
        densities.append(round(ratio, 5))
        if ratio < RENDER_BLANK_INK_RATIO:
            blank_pages.append(idx)
        elif ratio < RENDER_LOW_INK_RATIO:
            low_ink_pages.append(idx)
    result.metrics["rendered_ink_ratios"] = densities
    result.metrics["rendered_blank_pages"] = blank_pages
    result.metrics["rendered_low_ink_pages"] = low_ink_pages
    result.metrics["rendered_unreadable_pages"] = unreadable_pages
    if blank_pages:
        result.issues.append(
            f"Rendered {result.kind.upper()} has visually blank {page_label}(s): "
            + ", ".join(str(page) for page in blank_pages[:10])
        )
    if len(low_ink_pages) >= 3:
        result.issues.append(
            f"Rendered {result.kind.upper()} has too many visually sparse {page_label}(s): "
            + ", ".join(str(page) for page in low_ink_pages[:10])
        )


def _png_ink_ratio(path: Path) -> float | None:
    try:
        data = path.read_bytes()
        if not data.startswith(b"\x89PNG\r\n\x1a\n"):
            return None
        pos = 8
        width = height = color_type = bit_depth = None
        idat = bytearray()
        while pos + 8 <= len(data):
            length = struct.unpack(">I", data[pos : pos + 4])[0]
            chunk_type = data[pos + 4 : pos + 8]
            chunk = data[pos + 8 : pos + 8 + length]
            pos += 12 + length
            if chunk_type == b"IHDR":
                width, height, bit_depth, color_type, _compression, _filter, interlace = struct.unpack(
                    ">IIBBBBB", chunk
                )
                if interlace != 0 or bit_depth != 8 or color_type not in {0, 2, 6}:
                    return None
            elif chunk_type == b"IDAT":
                idat.extend(chunk)
            elif chunk_type == b"IEND":
                break
        if not width or not height or color_type is None:
            return None
        channels = {0: 1, 2: 3, 6: 4}[color_type]
        stride = width * channels
        raw = zlib.decompress(bytes(idat))
        rows: list[bytearray] = []
        offset = 0
        prev = bytearray(stride)
        for _row_idx in range(height):
            filter_type = raw[offset]
            offset += 1
            row = bytearray(raw[offset : offset + stride])
            offset += stride
            _unfilter_png_row(row, prev, filter_type, channels)
            rows.append(row)
            prev = row
        non_white = 0
        total = width * height
        for row in rows:
            for pixel in range(width):
                base = pixel * channels
                if color_type == 0:
                    rgb = (row[base], row[base], row[base])
                    alpha = 255
                else:
                    rgb = (row[base], row[base + 1], row[base + 2])
                    alpha = row[base + 3] if channels == 4 else 255
                if alpha > 10 and min(255 - channel for channel in rgb) > 18:
                    non_white += 1
        return non_white / total if total else None
    except Exception:
        return None


def _unfilter_png_row(row: bytearray, prev: bytearray, filter_type: int, bpp: int) -> None:
    for idx, value in enumerate(row):
        left = row[idx - bpp] if idx >= bpp else 0
        up = prev[idx] if prev else 0
        up_left = prev[idx - bpp] if prev and idx >= bpp else 0
        if filter_type == 0:
            continue
        if filter_type == 1:
            row[idx] = (value + left) & 0xFF
        elif filter_type == 2:
            row[idx] = (value + up) & 0xFF
        elif filter_type == 3:
            row[idx] = (value + ((left + up) // 2)) & 0xFF
        elif filter_type == 4:
            row[idx] = (value + _paeth(left, up, up_left)) & 0xFF
        else:
            raise ValueError(f"unsupported PNG filter type: {filter_type}")


def _paeth(left: int, up: int, up_left: int) -> int:
    estimate = left + up - up_left
    pa = abs(estimate - left)
    pb = abs(estimate - up)
    pc = abs(estimate - up_left)
    if pa <= pb and pa <= pc:
        return left
    if pb <= pc:
        return up
    return up_left


def _write_render_index(results: list[ArtifactQaResult], out_dir: Path) -> Path | None:
    image_entries: list[tuple[str, str, Path]] = []
    for result in results:
        for rendered_file in result.rendered_files:
            rendered_path = Path(rendered_file)
            if rendered_path.suffix.lower() != ".png":
                continue
            try:
                rel = rendered_path.relative_to(out_dir)
            except ValueError:
                rel = rendered_path
            image_entries.append((result.kind.upper(), rendered_path.stem, rel))
    if not image_entries:
        return None

    cards = []
    for kind, label, rel in image_entries:
        href = html.escape(rel.as_posix())
        cards.append(
            "\n".join(
                [
                    '<article class="card">',
                    f'  <div class="meta">{html.escape(kind)} · {html.escape(label)}</div>',
                    f'  <a href="{href}"><img src="{href}" alt="{html.escape(label)}" /></a>',
                    "</article>",
                ]
            )
        )

    page = f"""<!doctype html>
<html lang="en">
<head>
  <meta charset="utf-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1" />
  <title>Smart Report Artifact QA</title>
  <style>
    :root {{
      color-scheme: light;
      --ink: #111827;
      --muted: #6b7280;
      --line: #d7dce5;
      --paper: #f6f7f9;
      --panel: #ffffff;
    }}
    * {{ box-sizing: border-box; }}
    body {{
      margin: 0;
      background: var(--paper);
      color: var(--ink);
      font: 14px/1.45 Arial, Helvetica, sans-serif;
    }}
    header {{
      padding: 24px 28px 18px;
      border-bottom: 1px solid var(--line);
      background: var(--panel);
    }}
    h1 {{
      margin: 0 0 6px;
      font-size: 22px;
      line-height: 1.2;
      letter-spacing: 0;
    }}
    p {{ margin: 0; color: var(--muted); }}
    .rubric {{
      display: grid;
      grid-template-columns: repeat(auto-fit, minmax(260px, 1fr));
      gap: 8px 18px;
      margin-top: 16px;
      padding: 14px 0 0;
      border-top: 1px solid var(--line);
      color: var(--ink);
    }}
    .rubric div {{
      display: flex;
      gap: 8px;
      align-items: flex-start;
      color: var(--muted);
      font-size: 13px;
    }}
    .rubric b {{
      color: var(--ink);
      font-weight: 700;
    }}
    main {{
      display: grid;
      grid-template-columns: repeat(auto-fit, minmax(280px, 1fr));
      gap: 18px;
      padding: 22px;
    }}
    .card {{
      background: var(--panel);
      border: 1px solid var(--line);
      border-radius: 8px;
      overflow: hidden;
      box-shadow: 0 1px 2px rgba(15, 23, 42, 0.08);
    }}
    .meta {{
      padding: 10px 12px;
      border-bottom: 1px solid var(--line);
      color: var(--muted);
      font-size: 12px;
      text-transform: uppercase;
      letter-spacing: .08em;
    }}
    img {{
      display: block;
      width: 100%;
      height: auto;
      background: white;
    }}
  </style>
</head>
<body>
  <header>
    <h1>Smart Report Artifact QA</h1>
    <p>Rendered DOCX/PDF pages and PPTX slides for fast human visual review.</p>
    <section class="rubric" aria-label="Manual visual review rubric">
      {"".join(f"<div><b>{idx}.</b><span>{html.escape(item)}</span></div>" for idx, item in enumerate(VISUAL_REVIEW_RUBRIC, start=1))}
    </section>
  </header>
  <main>
    {"".join(cards)}
  </main>
</body>
</html>
"""
    index_path = out_dir / "index.html"
    index_path.write_text(page, encoding="utf-8")
    return index_path


def _add_common_content_issues(result: ArtifactQaResult, text: str) -> None:
    for marker in INTERNAL_MARKERS:
        if marker in text:
            result.issues.append(f"Internal marker leaked into {result.kind.upper()}: {marker}")
    if "\ufffd" in text:
        result.issues.append(f"{result.kind.upper()} contains replacement-character mojibake.")


def _require_metric(result: ArtifactQaResult, key: str, minimum: int, message: str) -> None:
    value = int(result.metrics.get(key) or 0)
    if value < minimum:
        result.issues.append(f"{message} Observed {value}, expected at least {minimum}.")


def _overall_status(results: list[ArtifactQaResult]) -> str:
    if not results:
        return "failed"
    if any(item.structural_status == "failed" for item in results):
        return "failed"
    if any(item.render_status == "failed" for item in results):
        return "failed"
    if any(item.render_status == "blocked" for item in results):
        return "blocked"
    return "passed"


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="QA premium Smart Report DOCX/PDF/PPTX artifacts.")
    parser.add_argument("--docx", type=Path, help="Path to premium DOCX report.")
    parser.add_argument("--pdf", type=Path, help="Path to premium publication PDF report.")
    parser.add_argument("--pptx", type=Path, help="Path to premium PPTX deck.")
    parser.add_argument("--out-dir", type=Path, default=Path("tmp/premium_artifact_qa"))
    parser.add_argument("--json", type=Path, help="Optional JSON output path.")
    parser.add_argument("--no-render", action="store_true", help="Skip LibreOffice/Poppler rendering.")
    parser.add_argument("--strict", action="store_true", help="Return non-zero for blocked visual render QA.")
    args = parser.parse_args(argv)

    report = run_qa(
        docx_path=args.docx,
        pdf_path=args.pdf,
        pptx_path=args.pptx,
        out_dir=args.out_dir,
        render=not args.no_render,
    )
    payload = json.dumps(report, ensure_ascii=False, indent=2)
    print(payload)
    if args.json:
        args.json.parent.mkdir(parents=True, exist_ok=True)
        args.json.write_text(payload + "\n", encoding="utf-8")
    if report["status"] == "failed":
        return 1
    if args.strict and report["status"] == "blocked":
        return 2
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
