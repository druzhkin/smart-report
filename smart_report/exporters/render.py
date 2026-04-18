"""Format renderers for the uniform report-dict produced by v4_to_report.

Seven entry points correspond to the seven formats in the existing
ExportDropdown:

    write_md, write_json, write_onepager_html   (zero external deps)
    write_docx, write_pptx                      (python-docx, python-pptx)
    write_gamma_pptx_stub, write_gamma_pdf_stub (placeholder files + json
        companion — Gamma API integration is future work)

The Gamma stubs write a small .json manifest the UI can distinguish from a
real export, plus a plaintext placeholder so the download stream doesn't
fail when no GAMMA_API_KEY is configured. This matches spec §4 Track B
bullet 8's "Gamma can be mocked if no API key" line.
"""

from __future__ import annotations

import html
import json
from pathlib import Path
from typing import Any

REPORT_DICT = dict[str, Any]


# ---------- markdown ----------


def render_markdown(rd: REPORT_DICT) -> str:
    """Return the full report as a markdown string."""
    lines: list[str] = []
    title = rd.get("title") or "Smart Report v4"
    lines.append(f"# {title}\n")

    q = rd.get("question") or ""
    if q:
        lines.append(f"**Вопрос:** {q}\n")

    es = rd.get("executive_summary") or {}
    lines.append("\n## Executive Summary\n")

    main_answer = (es.get("main_answer") or "").strip()
    if main_answer:
        lines.append(main_answer + "\n")

    ranking = es.get("ranking")
    if ranking:
        lines.append(f"\n**Ранжирование:** {ranking}\n")

    top = es.get("top_findings") or []
    if top:
        lines.append("\n### Ключевые находки\n")
        for tf in top:
            lines.append(f"- {tf}")
        lines.append("")

    kns = es.get("key_numbers") or []
    if kns:
        lines.append("\n### Ключевые цифры\n")
        for kn in kns:
            src = kn.get("source_url") or ""
            src_md = f" ([источник]({src}))" if src else ""
            subj = kn.get("subject") or ""
            subj_tail = f" — {subj}" if subj else ""
            lines.append(
                f"- **{kn.get('value','')}** · {kn.get('metric','')}{subj_tail}{src_md}"
            )
        lines.append("")

    cn = (es.get("confidence_note") or "").strip()
    if cn:
        lines.append(f"\n**Оценка уверенности:** {cn}\n")
    wma = (es.get("what_meta_adds") or "").strip()
    if wma:
        lines.append(f"\n**Что даёт мета-анализ:** {wma}\n")

    sections = rd.get("sections") or []
    for sec in sections:
        lines.append(f"\n## {sec.get('heading', '')}\n")
        lines.append(sec.get("body_markdown", "").strip() + "\n")

    sources = rd.get("sources") or []
    if sources:
        lines.append("\n## Источники\n")
        for s in sources:
            url = s.get("url") or ""
            link = f" <{url}>" if url else ""
            rel = s.get("reliability") or ""
            rel_tail = f" · _{rel}_" if rel else ""
            tool = s.get("tool") or ""
            tool_tail = f" · {tool}" if tool else ""
            lines.append(f"- **{s.get('title','(без названия)')}**{link}{tool_tail}{rel_tail}")
        lines.append("")

    meta = rd.get("metadata") or {}
    if meta:
        lines.append("\n---\n_Метаданные:_ " + ", ".join(f"{k}={v}" for k, v in meta.items()))

    return "\n".join(lines) + "\n"


def write_md(path: Path, rd: REPORT_DICT) -> Path:
    path = Path(path)
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(render_markdown(rd), encoding="utf-8")
    return path


def write_json(path: Path, rd: REPORT_DICT) -> Path:
    path = Path(path)
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(
        json.dumps(rd, ensure_ascii=False, indent=2, default=str), encoding="utf-8"
    )
    return path


# ---------- one-pager HTML ----------


def write_onepager_html(path: Path, rd: REPORT_DICT) -> Path:
    path = Path(path)
    path.parent.mkdir(parents=True, exist_ok=True)

    es = rd.get("executive_summary") or {}
    title = html.escape(rd.get("title") or "Smart Report v4")
    q = html.escape(rd.get("question") or "")
    main = html.escape(es.get("main_answer") or "").replace("\n", "<br/>")
    ranking = es.get("ranking")
    top = es.get("top_findings") or []
    kns = es.get("key_numbers") or []
    sources = rd.get("sources") or []

    parts: list[str] = [
        "<!DOCTYPE html>",
        "<html lang='ru'><head>",
        "<meta charset='utf-8'/>",
        f"<title>{title}</title>",
        "<style>",
        "body{font-family:-apple-system,Segoe UI,Inter,sans-serif;max-width:820px;"
        "margin:40px auto;padding:0 24px;color:#111;line-height:1.55;}",
        "h1{font-size:22px;margin-bottom:4px;} h2{font-size:16px;margin-top:24px;}",
        ".q{color:#555;font-size:13px;margin-bottom:20px;}",
        ".main{background:#f6f7f9;border-left:4px solid #2563eb;padding:12px 16px;"
        "border-radius:6px;font-size:15px;margin:16px 0;}",
        ".kn{display:inline-block;background:#eef2ff;border:1px solid #c7d2fe;"
        "padding:2px 8px;border-radius:999px;margin:2px;font-size:12px;}",
        "ul{padding-left:20px;} li{margin:4px 0;}",
        "a{color:#2563eb;text-decoration:none;} a:hover{text-decoration:underline;}",
        ".meta{color:#888;font-size:12px;margin-top:30px;border-top:1px solid #eee;"
        "padding-top:10px;}",
        "</style></head><body>",
        f"<h1>{title}</h1>",
        f"<div class='q'>{q}</div>",
        f"<div class='main'>{main}</div>",
    ]
    if ranking:
        parts.append(f"<div><strong>Ранжирование:</strong> {html.escape(ranking)}</div>")
    if top:
        parts.append("<h2>Ключевые находки</h2><ul>")
        for tf in top:
            parts.append(f"<li>{html.escape(str(tf))}</li>")
        parts.append("</ul>")
    if kns:
        parts.append("<h2>Ключевые цифры</h2><div>")
        for kn in kns:
            v = html.escape(str(kn.get("value", "")))
            m = html.escape(str(kn.get("metric", "")))
            parts.append(f"<span class='kn'>{v} · {m}</span>")
        parts.append("</div>")
    if sources:
        parts.append("<h2>Источники</h2><ul>")
        for s in sources:
            t = html.escape(str(s.get("title") or "(без названия)"))
            u = s.get("url") or ""
            if u:
                parts.append(f"<li><a href='{html.escape(u)}'>{t}</a></li>")
            else:
                parts.append(f"<li>{t}</li>")
        parts.append("</ul>")
    meta = rd.get("metadata") or {}
    if meta:
        meta_str = ", ".join(
            f"{html.escape(str(k))}={html.escape(str(v))}" for k, v in meta.items()
        )
        parts.append(f"<div class='meta'>{meta_str}</div>")
    parts.append("</body></html>")
    path.write_text("\n".join(parts), encoding="utf-8")
    return path


# ---------- docx ----------


def write_docx(path: Path, rd: REPORT_DICT) -> Path:
    from docx import Document  # lazy import so md/json work without python-docx

    path = Path(path)
    path.parent.mkdir(parents=True, exist_ok=True)

    doc = Document()
    doc.add_heading(rd.get("title") or "Smart Report v4", level=0)

    q = rd.get("question")
    if q:
        p = doc.add_paragraph()
        p.add_run("Вопрос: ").bold = True
        p.add_run(q)

    es = rd.get("executive_summary") or {}
    doc.add_heading("Executive Summary", level=1)
    main = es.get("main_answer") or ""
    if main:
        doc.add_paragraph(main)

    if es.get("ranking"):
        p = doc.add_paragraph()
        p.add_run("Ранжирование: ").bold = True
        p.add_run(es["ranking"])

    if es.get("top_findings"):
        doc.add_heading("Ключевые находки", level=2)
        for tf in es["top_findings"]:
            doc.add_paragraph(tf, style="List Bullet")

    if es.get("key_numbers"):
        doc.add_heading("Ключевые цифры", level=2)
        for kn in es["key_numbers"]:
            subj = f" — {kn.get('subject','')}" if kn.get("subject") else ""
            src = f" ({kn.get('source_url','')})" if kn.get("source_url") else ""
            doc.add_paragraph(
                f"{kn.get('value','')} · {kn.get('metric','')}{subj}{src}",
                style="List Bullet",
            )
    if es.get("confidence_note"):
        p = doc.add_paragraph()
        p.add_run("Оценка уверенности: ").bold = True
        p.add_run(es["confidence_note"])
    if es.get("what_meta_adds"):
        p = doc.add_paragraph()
        p.add_run("Что даёт мета-анализ: ").bold = True
        p.add_run(es["what_meta_adds"])

    for sec in rd.get("sections") or []:
        doc.add_heading(sec.get("heading", ""), level=1)
        for para in (sec.get("body_markdown") or "").split("\n\n"):
            if para.strip():
                doc.add_paragraph(para.strip())

    sources = rd.get("sources") or []
    if sources:
        doc.add_heading("Источники", level=1)
        for s in sources:
            line = s.get("title") or "(без названия)"
            if s.get("url"):
                line += f" — {s['url']}"
            if s.get("reliability"):
                line += f" [{s['reliability']}]"
            doc.add_paragraph(line, style="List Bullet")

    doc.save(str(path))
    return path


# ---------- pptx ----------


def write_pptx(path: Path, rd: REPORT_DICT) -> Path:
    from pptx import Presentation  # lazy import
    from pptx.util import Inches, Pt

    path = Path(path)
    path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    # Title slide.
    blank = prs.slide_layouts[5]  # title only
    slide = prs.slides.add_slide(blank)
    slide.shapes.title.text = rd.get("title") or "Smart Report v4"

    es = rd.get("executive_summary") or {}

    # Executive summary slide.
    slide = prs.slides.add_slide(blank)
    slide.shapes.title.text = "Executive Summary"
    tx = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.3), Inches(9), Inches(5.5)
    ).text_frame
    tx.word_wrap = True
    _add_para(tx, es.get("main_answer") or "", size=14)
    if es.get("ranking"):
        _add_para(tx, f"Ранжирование: {es['ranking']}", size=13, bold=True)
    for tf in es.get("top_findings") or []:
        _add_para(tx, f"• {tf}", size=12)

    # Key numbers slide.
    if es.get("key_numbers"):
        slide = prs.slides.add_slide(blank)
        slide.shapes.title.text = "Ключевые цифры"
        tx = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.3), Inches(9), Inches(5.5)
        ).text_frame
        tx.word_wrap = True
        for kn in es["key_numbers"]:
            line = f"{kn.get('value','')} · {kn.get('metric','')}"
            if kn.get("subject"):
                line += f" — {kn['subject']}"
            _add_para(tx, line, size=14)

    # Per-section slides.
    for sec in rd.get("sections") or []:
        slide = prs.slides.add_slide(blank)
        slide.shapes.title.text = sec.get("heading", "")
        tx = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.3), Inches(9), Inches(5.5)
        ).text_frame
        tx.word_wrap = True
        body = (sec.get("body_markdown") or "").strip()
        # PPTX doesn't render markdown — just chunk long body so the box
        # isn't one wall. Cap at ~2000 chars so a single slide stays usable.
        chunks = [body[i : i + 1000] for i in range(0, min(len(body), 2000), 1000)]
        for chunk in chunks:
            _add_para(tx, chunk, size=11)
        if len(body) > 2000:
            _add_para(tx, "[...] (truncated — see docx/md for full text)", size=10)

    # Sources slide.
    sources = rd.get("sources") or []
    if sources:
        slide = prs.slides.add_slide(blank)
        slide.shapes.title.text = "Источники"
        tx = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.3), Inches(9), Inches(5.5)
        ).text_frame
        tx.word_wrap = True
        for s in sources[:15]:
            line = s.get("title") or "(без названия)"
            if s.get("url"):
                line += f" — {s['url']}"
            _add_para(tx, line, size=10)

    prs.save(str(path))
    return path


def _add_para(tf, text: str, *, size: int = 12, bold: bool = False) -> None:
    from pptx.util import Pt

    p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
    p.text = text
    for run in p.runs:
        run.font.size = Pt(size)
        if bold:
            run.font.bold = True


# ---------- Gamma stubs ----------
# Until a GAMMA_API_KEY is wired, the Gamma formats write a small JSON manifest
# next to a plaintext placeholder. The UI's ExportDropdown receives a file, the
# file tells the user Gamma was not called. Swapping in real Gamma is a later
# task (see v2's export_gamma.py for the shape).


def write_gamma_pptx_stub(path: Path, rd: REPORT_DICT) -> Path:
    return _gamma_stub(path, rd, kind="gamma-pptx")


def write_gamma_pdf_stub(path: Path, rd: REPORT_DICT) -> Path:
    return _gamma_stub(path, rd, kind="gamma-pdf")


def _gamma_stub(path: Path, rd: REPORT_DICT, *, kind: str) -> Path:
    path = Path(path)
    path.parent.mkdir(parents=True, exist_ok=True)
    payload = {
        "stub": True,
        "kind": kind,
        "reason": (
            "GAMMA_API_KEY is not configured in this environment. "
            "Falling back to a stub file so the download stream does not fail. "
            "For real Gamma output, wire GAMMA_API_KEY and replace this handler."
        ),
        "title": rd.get("title"),
        "session_id": rd.get("session_id"),
        "main_answer": (rd.get("executive_summary") or {}).get("main_answer"),
    }
    path.write_text(
        json.dumps(payload, ensure_ascii=False, indent=2, default=str),
        encoding="utf-8",
    )
    return path
