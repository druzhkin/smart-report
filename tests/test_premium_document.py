from __future__ import annotations

import pytest

from smart_report.exporters import v4_to_report_dict
from smart_report.exporters.premium import (
    CarboneRenderError,
    assemble_premium_report_document,
    render_premium_carbone_pdf,
    render_premium_docx,
    render_premium_pdf,
    render_premium_pptx,
    to_carbone_data,
)
from smart_report.models import (
    AnalysisOutput,
    ChartSpec,
    Conflict,
    ConsensusClaim,
    ExecutiveSummaryV4,
    FinalReport,
    Gap,
    KeyNumberHighlight,
    NumericFact,
    Source,
    SourceRef,
)


def test_render_premium_pdf_opens_and_meets_publication_shape(tmp_path):
    pytest.importorskip("reportlab")
    pypdf = pytest.importorskip("pypdf")

    document = assemble_premium_report_document(
        _report(),
        analysis=_analysis(),
        premium_readiness={
            "ready": True,
            "score": 91,
            "issues": [],
            "strengths": ["Publication-grade layout requirements are present."],
        },
    )
    out = render_premium_pdf(document, tmp_path / "premium_report.pdf")

    assert out.exists()
    assert out.stat().st_size > 1000
    reader = pypdf.PdfReader(str(out))
    assert len(reader.pages) >= document.plan.deliverables.report_min_pages
    first_page_text = reader.pages[0].extract_text() or ""
    assert "SMART REPORT" in first_page_text
    assert "Publication-grade PDF" in first_page_text
    publication_text = "\n".join(page.extract_text() or "" for page in reader.pages[:8])
    assert "EXHIBIT" in publication_text
    assert "Индексированный числовой сигнал" in publication_text
    assert "Надежность источников" in publication_text
    assert "Interpretation" not in publication_text
    assert "What it means" not in publication_text
    assert "RANKING BAR" not in publication_text
    assert "Source reliability mix" not in publication_text
    assert "Source:" in publication_text


def test_carbone_data_flattens_premium_document_for_template():
    document = assemble_premium_report_document(_report(), analysis=_analysis())

    data = to_carbone_data(document)

    assert data["title"] == document.title
    assert data["sourceCount"] == 2
    assert data["numericFactCount"] == 1
    assert len(data["sections"]) == len(document.sections)
    assert len(data["appendices"]) == len(document.appendices)
    assert len(data["visualPages"]) == len(document.pages)
    assert any(page["chartFrameClass"] == "" for page in data["visualPages"])
    assert any(page["kpiClass"] == "" for page in data["visualPages"])
    assert len(data["exhibits"]) >= document.plan.publication.min_exhibit_pages
    assert "<table>" in data["sections"][0]["blocksHtml"] or "<p>" in data["sections"][0]["blocksHtml"]


def test_carbone_renderer_requires_env_token(tmp_path, monkeypatch):
    monkeypatch.delenv("CARBONE_API_KEY", raising=False)
    monkeypatch.delenv("CARBONE_TOKEN", raising=False)
    document = assemble_premium_report_document(_report(), analysis=_analysis())

    with pytest.raises(CarboneRenderError, match="CARBONE_API_KEY"):
        render_premium_carbone_pdf(document, tmp_path / "premium.pdf")


def test_carbone_renderer_posts_inline_template_and_writes_pdf(tmp_path, monkeypatch):
    document = assemble_premium_report_document(_report(), analysis=_analysis())
    calls = []

    class _FakeClient:
        def __init__(self, *args, **kwargs):
            pass

        def __enter__(self):
            return self

        def __exit__(self, *args):
            return None

        def post(self, url, *, headers, json):
            calls.append({"url": url, "headers": headers, "json": json})
            import httpx

            return httpx.Response(
                200,
                content=b"%PDF-1.4\n% mock\n",
                request=httpx.Request("POST", url),
                headers={"content-type": "application/pdf"},
            )

    monkeypatch.setattr("smart_report.exporters.premium.carbone.httpx.Client", _FakeClient)

    out = render_premium_carbone_pdf(
        document,
        tmp_path / "premium.pdf",
        api_token="test-token",
        api_url="https://api.carbone.test",
    )

    assert out.read_bytes().startswith(b"%PDF")
    assert calls[0]["url"] == "https://api.carbone.test/render/template?download=true"
    assert calls[0]["headers"]["Authorization"] == "Bearer test-token"
    assert calls[0]["headers"]["carbone-version"] == "5"
    assert calls[0]["json"]["convertTo"] == "pdf"
    assert calls[0]["json"]["converter"] == "C"
    assert calls[0]["json"]["template"]
    assert calls[0]["json"]["data"]["title"] == document.title


def _report() -> FinalReport:
    return FinalReport(
        session_id="premium-doc",
        question="Forecast a market with scenario and risk recommendations",
        executive_summary=ExecutiveSummaryV4(
            main_answer="Base case is moderate growth with material downside triggers.",
            top_findings=["Demand is rate-sensitive.", "Supply remains constrained."],
            confidence_note="Confidence is medium because transaction data is partial.",
            what_meta_adds="Consensus, conflicts, and gaps are separated.",
        ),
        main_synthesis="The market baseline combines price, demand, supply, and policy drivers.",
        consensus_section="Sources agree on the direction but differ on magnitude.",
        conflicts_section="Sources disagree on timing.",
        gaps_filled_section="Project-level microdata remains unavailable.",
        all_sources=[
            Source(title="Official source", url="https://example.com/official", reliability="high"),
            Source(title="Industry source", url="https://example.com/industry", reliability="medium"),
        ],
        charts=[
            ChartSpec(
                chart_type="bar",
                title="Market driver ranking",
                data={
                    "points": [
                        {"label": "Demand", "value": 10},
                        {"label": "Supply", "value": 7},
                        {"label": "Rates", "value": 4},
                    ]
                },
                caption="Synthetic fixture chart for premium storyboard tests.",
            )
        ],
        key_numbers_highlight=[
            KeyNumberHighlight(
                value="10%",
                label="Market growth signal",
                source_ref="https://example.com/official",
                importance="headline",
            )
        ],
    )


def _analysis() -> AnalysisOutput:
    facts = [
        NumericFact(
            fact_id="f1",
            value="10%",
            metric="growth",
            subject="market",
            relevance_to_question="high",
            sources=[SourceRef(url="https://example.com/official", title="Official")],
        )
    ]
    return AnalysisOutput(
        consensus=[
            ConsensusClaim(claim="Demand is sensitive to financing conditions.", confidence="high")
        ],
        conflicts=[
            Conflict(
                topic="Growth range",
                source_a="A",
                claim_a="High growth",
                source_b="B",
                claim_b="Moderate growth",
                importance="material",
            )
        ],
        gaps=[
            Gap(
                topic="Transaction microdata",
                why_critical="Needed for exact pricing",
                what_to_find="Closed transaction dataset",
            )
        ],
        all_numeric_facts=facts,
        high_relevance_facts=facts,
    )


def test_assemble_premium_document_uses_existing_analysis_layers():
    document = assemble_premium_report_document(_report(), analysis=_analysis())

    assert document.plan.deliverables.report_min_pages >= 20
    assert document.plan.deliverables.deck_min_slides >= 10
    assert document.source_count == 2
    assert document.numeric_fact_count == 1
    assert len(document.sections) >= 7
    assert len(document.appendices) == 3
    assert len(document.deck_slides) >= 10
    assert len(document.pages) >= 8
    assert all("Позиция автора" not in page.thesis for page in document.pages)

    visual_pages = [page for page in document.pages if page.visual and page.visual.visual_type != "none"]
    visual_types = {page.visual.visual_type for page in visual_pages if page.visual}
    assert len(visual_pages) / len(document.pages) >= 0.6
    assert {"hero_kpi_strip", "ranking_bar", "risk_heatmap", "evidence_quality", "source_table"} <= visual_types
    assert all(page.thesis for page in document.pages)

    block_titles = {
        block.title
        for section in [*document.sections, *document.appendices]
        for block in section.blocks
    }
    assert "Реестр числовых доказательств" in block_titles
    assert "Противоречия и расхождения" in block_titles
    assert "Реестр рисков" in block_titles
    assert "Матрица решений" in block_titles


def test_assemble_premium_document_does_not_mutate_legacy_report():
    report = _report()
    before = v4_to_report_dict(report)
    _ = assemble_premium_report_document(report, analysis=_analysis())
    after = v4_to_report_dict(report)

    assert before == after


def test_render_premium_docx_opens_and_contains_report_structure(tmp_path):
    pytest.importorskip("docx")
    from docx import Document

    document = assemble_premium_report_document(
        _report(),
        analysis=_analysis(),
        premium_readiness={
            "ready": False,
            "score": 61,
            "issues": [
                {
                    "code": "premium_too_few_authoritative_sources",
                    "severity": "critical",
                    "message": "Authoritative source threshold is not met.",
                    "recommendation": "Add primary sources before delivery.",
                }
            ],
            "strengths": ["Consensus layer is present."],
        },
    )
    out = render_premium_docx(document, tmp_path / "premium.docx")

    assert out.exists()
    assert out.stat().st_size > 1000
    loaded = Document(out)
    table_text = "\n".join(
        cell.text
        for table in loaded.tables
        for row in table.rows
        for cell in row.cells
    )
    text = "\n".join([*(p.text for p in loaded.paragraphs), table_text])
    assert "SMART REPORT | ПРЕМИАЛЬНЫЙ АНАЛИТИЧЕСКИЙ ОТЧЁТ" in text
    assert "Резюме для решения" in text
    assert "Короткий ответ" in text
    assert "Карта доказательств" in text
    assert "НЕ ГОТОВ К ПЛАТНОЙ ВЫДАЧЕ КЛИЕНТУ" not in text
    assert "Как читать отчёт" in text
    assert "Реестр числовых доказательств" in text
    assert len(loaded.tables) >= 4


def test_render_premium_docx_can_include_internal_audit_when_requested(tmp_path):
    pytest.importorskip("docx")
    from docx import Document

    document = assemble_premium_report_document(
        _report(),
        analysis=_analysis(),
        premium_readiness={
            "ready": False,
            "score": 61,
            "issues": [
                {
                    "code": "premium_too_few_authoritative_sources",
                    "severity": "critical",
                    "message": "Authoritative source threshold is not met.",
                    "recommendation": "Add primary sources before delivery.",
                }
            ],
            "strengths": ["Consensus layer is present."],
        },
    )
    out = render_premium_docx(
        document,
        tmp_path / "premium_internal.docx",
        include_internal_audit=True,
    )

    loaded = Document(out)
    table_text = "\n".join(
        cell.text
        for table in loaded.tables
        for row in table.rows
        for cell in row.cells
    )
    text = "\n".join([*(p.text for p in loaded.paragraphs), table_text])
    assert "Гейт готовности к платной выдаче" in text
    assert "НЕ ГОТОВ К ПЛАТНОЙ ВЫДАЧЕ КЛИЕНТУ" in text


def test_render_premium_pptx_opens_and_contains_deck_structure(tmp_path):
    pytest.importorskip("pptx")
    from pptx import Presentation

    document = assemble_premium_report_document(
        _report(),
        analysis=_analysis(),
        premium_readiness={
            "ready": False,
            "score": 61,
            "issues": [
                {
                    "code": "premium_too_few_authoritative_sources",
                    "severity": "critical",
                    "message": "Authoritative source threshold is not met.",
                    "recommendation": "Add primary sources before delivery.",
                }
            ],
            "strengths": ["Consensus layer is present."],
        },
    )
    out = render_premium_pptx(document, tmp_path / "premium_deck.pptx")

    assert out.exists()
    assert out.stat().st_size > 1000
    deck = Presentation(str(out))
    assert len(deck.slides) >= 10
    slide_text = "\n".join(
        shape.text
        for slide in deck.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )
    assert "Короткий ответ" in slide_text
    assert "Позиция и ограничения" in slide_text
    assert "Готовность к платной выдаче" not in slide_text
    assert "НЕ ГОТОВ" not in slide_text


def test_render_premium_pptx_can_include_internal_audit_when_requested(tmp_path):
    pytest.importorskip("pptx")
    from pptx import Presentation

    document = assemble_premium_report_document(
        _report(),
        analysis=_analysis(),
        premium_readiness={
            "ready": False,
            "score": 61,
            "issues": [
                {
                    "code": "premium_too_few_authoritative_sources",
                    "severity": "critical",
                    "message": "Authoritative source threshold is not met.",
                    "recommendation": "Add primary sources before delivery.",
                }
            ],
            "strengths": ["Consensus layer is present."],
        },
    )
    out = render_premium_pptx(
        document,
        tmp_path / "premium_internal_deck.pptx",
        include_internal_audit=True,
    )

    deck = Presentation(str(out))
    slide_text = "\n".join(
        shape.text
        for slide in deck.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )
    assert "Готовность к платной выдаче" in slide_text
    assert "НЕ ГОТОВ" in slide_text
